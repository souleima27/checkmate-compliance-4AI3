import requests
import os
import json
from pptx import Presentation
from reportlab.pdfgen import canvas

BASE_URL = "http://localhost:8000/api"

def create_dummy_pptx(filename):
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Hello, World!"
    subtitle.text = "This is a test presentation."
    prs.save(filename)

def create_dummy_pdf(filename):
    c = canvas.Canvas(filename)
    c.drawString(100, 750, "Hello, World!")
    c.drawString(100, 730, "This is a test PDF document.")
    c.save()

def test_audit():
    print("Testing /api/audit...")
    
    # Create dummy files
    try:
        create_dummy_pptx("test.pptx")
        create_dummy_pdf("test.pdf")
    except ImportError:
        print("Required libraries (python-pptx, reportlab) not found. Skipping file creation.")
        return

    files = {
        "pptx_file": ("test.pptx", open("test.pptx", "rb"), "application/vnd.openxmlformats-officedocument.presentationml.presentation"),
        "pdf_file": ("test.pdf", open("test.pdf", "rb"), "application/pdf")
    }
    
    try:
        response = requests.post(f"{BASE_URL}/audit", files=files)
        print(f"Status Code: {response.status_code}")
        if response.status_code == 200:
            print("Response JSON keys:", response.json().keys())
            print("Audit successful!")
        else:
            print("Audit failed:", response.text)
    except Exception as e:
        print(f"Request failed: {e}")
    finally:
        # Cleanup
        if os.path.exists("test.pptx"):
            try:
                files["pptx_file"][1].close()
                os.remove("test.pptx")
            except: pass
        if os.path.exists("test.pdf"):
            try:
                files["pdf_file"][1].close()
                os.remove("test.pdf")
            except: pass

def test_chat():
    print("\nTesting /api/chat...")
    
    payload = {
        "question": "Quelle est l'estimation du capital requis ?"
    }
    
    try:
        response = requests.post(f"{BASE_URL}/chat", json=payload)
        print(f"Status Code: {response.status_code}")
        if response.status_code == 200:
            data = response.json()
            print("Answer:", data.get("answer")[:50] + "...")
            print("Sources present:", "sources" in data)
            print("Metrics present:", "metrics" in data)
        else:
            print("Chat failed:", response.text)
    except Exception as e:
        print(f"Request failed: {e}")

if __name__ == "__main__":
    test_audit()
    test_chat()
