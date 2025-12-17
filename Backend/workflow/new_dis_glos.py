"""
DISCLAIMER & GLOSSARY AGENT - VERSION INTELLIGENTE OPTIMISÉE
--------------------------------------------------------------------------------
Agent de conformité intelligent pour l'analyse des présentations PPTX marketing.
Version optimisée avec réduction des appels LLM et corrections des erreurs.

Fonctionnalités principales :
1. Analyse intelligente des disclaimers et de leur alignement
2. Détection contextuelle des sources
3. Vérification intelligente des glossaires
4. Extraction intelligente des caractéristiques du fonds
5. Comparaison avec la base de données des fonds enregistrés
6. Validation des pays de commercialisation

Auteur: Assistant IA
Date: 2025-12-12
Version: 7.0_optimisee
"""

import json
import os
import sys
import re
import time
import warnings
from typing import Dict, List, Any, Optional, Tuple, TypedDict
from datetime import datetime
import numpy as np
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from langgraph.graph import StateGraph, END
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import unicodedata

# Ajout du chemin pour les imports
current_dir = os.path.dirname(os.path.abspath(__file__))
root_dir = os.path.dirname(current_dir)
sys.path.append(root_dir)
sys.path.append(current_dir)

warnings.filterwarnings('ignore')

# Import des modules centralisés
try:
    from config.settings import *
    from llm_client import LlamaClient
    print("✅ Modules config.settings et llm_client importés avec succès")
except ImportError as e:
    print(f"❌ Erreur import modules: {e}")
    print("Veuillez vérifier que les fichiers config/settings.py et llm_client.py existent")
    sys.exit(1)

# ==================== CONFIGURATION ====================
CACHE_DIR = "workflow/caches"
FOND_REGISTRED_PATH = os.path.join(CACHE_DIR, "fond_registred.json")
DISCLAIMERS_PATH = os.path.join(CACHE_DIR, "disclaimers.json")
GLOSSAIRES_PATH = os.path.join(CACHE_DIR, "glossaires.json")

# ==================== IMPLEMENTATION LEVENSHTEIN FALLBACK ====================
try:
    import Levenshtein
    print("✅ Module Levenshtein importé avec succès")
except ImportError:
    print("⚠️  Module Levenshtein non trouvé, utilisation de l'implémentation de secours")
    
    class Levenshtein:
        """Implémentation de secours de Levenshtein"""
        @staticmethod
        def ratio(s1: str, s2: str) -> float:
            """Calcule la similarité entre deux chaînes (0-1)"""
            if not s1 or not s2:
                return 0.0
            
            s1, s2 = s1.lower().strip(), s2.lower().strip()
            
            # Si une chaîne est contenue dans l'autre
            if s1 in s2 or s2 in s1:
                return 0.9
            
            # Similarité basée sur les mots communs
            words1 = set(s1.split())
            words2 = set(s2.split())
            
            if not words1 or not words2:
                return 0.0
            
            intersection = len(words1.intersection(words2))
            union = len(words1.union(words2))
            
            return intersection / union if union > 0 else 0.0

# ==================== UTILITAIRES AVANCÉS DE TEXTE ====================
class TextUtils:
    """Utilitaires avancés pour le traitement du texte"""
    
    @staticmethod
    def normalize_text(text: str) -> str:
        """Normalise le texte pour la comparaison"""
        if not text:
            return ""
        
        # Convertir en minuscules
        text = text.lower()
        
        # Supprimer les accents
        text = unicodedata.normalize('NFKD', text).encode('ASCII', 'ignore').decode('ASCII')
        
        # Remplacer les caractères spéciaux
        text = re.sub(r'[^\w\s]', ' ', text)
        
        # Supprimer les espaces multiples
        text = re.sub(r'\s+', ' ', text).strip()
        
        return text
    
    @staticmethod
    def is_valid_isin(isin: str) -> bool:
        """Vérifie si une chaîne est un ISIN valide"""
        if not isin or len(isin) != 12:
            return False
        
        # Doit commencer par 2 lettres
        if not re.match(r'^[A-Z]{2}', isin):
            return False
        
        # Doit contenir au moins un chiffre
        if not any(c.isdigit() for c in isin):
            return False
        
        # Ne doit pas être uniquement des lettres
        if isin.isalpha():
            return False
        
        # Vérifier les faux positifs communs
        common_false_positives = [
            'PROMOTIONNEL', 'CONTRIBUTION', 'INFORMATION', 'PRESENTATION',
            'COMMERCIAL', 'PROFESSIONNEL', 'DISTRIBUTION', 'INVESTISSEMENT'
        ]
        if isin in common_false_positives:
            return False
        
        return True
    
    @staticmethod
    def find_similar_strings(search_list: List[str], target_list: List[str], threshold: float = 0.7) -> Dict[str, List[str]]:
        """Trouve les chaînes similaires entre deux listes"""
        similar_matches = {}
        
        for search_item in search_list:
            search_norm = TextUtils.normalize_text(search_item)
            matches = []
            
            for target_item in target_list:
                target_norm = TextUtils.normalize_text(target_item)
                
                # Score de similarité
                if search_norm == target_norm:
                    score = 1.0
                elif search_norm in target_norm or target_norm in search_norm:
                    score = 0.9
                else:
                    score = Levenshtein.ratio(search_norm, target_norm)
                
                if score >= threshold:
                    matches.append({
                        "original": target_item,
                        "normalized": target_norm,
                        "score": round(score, 3)
                    })
            
            if matches:
                # Trier par score décroissant
                matches.sort(key=lambda x: x["score"], reverse=True)
                similar_matches[search_item] = matches
        
        return similar_matches

# ==================== DÉTECTEUR DE PAYS INTELLIGENT ====================
class IntelligentCountryDetector:
    """Détecteur intelligent de pays avec validation améliorée"""
    
    # Liste complète des pays avec leurs noms en différentes langues
    COUNTRIES_DATABASE = {
        # Europe
        "Germany": ["Allemagne", "Deutschland", "Germania", "Alemania"],
        "France": ["France", "Frankreich", "Francia"],
        "United Kingdom": ["United Kingdom", "UK", "Royaume-Uni", "Großbritannien"],
        "Italy": ["Italy", "Italie", "Italien", "Italia"],
        "Spain": ["Spain", "Espagne", "Spanien", "España"],
        "Switzerland": ["Switzerland", "Suisse", "Schweiz", "Svizzera"],
        "Netherlands": ["Netherlands", "Pays-Bas", "Niederlande", "Países Bajos"],
        "Belgium": ["Belgium", "Belgique", "Belgien", "Bélgica"],
        "Luxembourg": ["Luxembourg", "Luxemburg", "Luxemburgo"],
        "Austria": ["Austria", "Autriche", "Österreich", "Austria"],
        "Portugal": ["Portugal", "Portugal", "Portugal", "Portugal"],
        "Ireland": ["Ireland", "Irlande", "Irland", "Irlanda"],
        "Sweden": ["Sweden", "Suède", "Schweden", "Suecia"],
        "Norway": ["Norway", "Norvège", "Norwegen", "Noruega"],
        "Denmark": ["Denmark", "Danemark", "Dänemark", "Dinamarca"],
        "Finland": ["Finland", "Finlande", "Finnland", "Finlandia"],
        "Poland": ["Poland", "Pologne", "Polen", "Polonia"],
        "Czech Republic": ["Czech Republic", "République Tchèque", "Tschechien", "República Checa"],
        "Hungary": ["Hungary", "Hongrie", "Ungarn", "Hungría"],
        "Slovakia": ["Slovakia", "Slovaquie", "Slowakei", "Eslovaquia"],
        "Slovenia": ["Slovenia", "Slovénie", "Slowenien", "Eslovenia"],
        "Croatia": ["Croatia", "Croatie", "Kroatien", "Croacia"],
        "Romania": ["Romania", "Roumanie", "Rumänien", "Rumanía"],
        "Bulgaria": ["Bulgaria", "Bulgarie", "Bulgarien", "Bulgaria"],
        "Greece": ["Greece", "Grèce", "Griechenland", "Grecia"],
        "Cyprus": ["Cyprus", "Chypre", "Zypern", "Chipre"],
        "Malta": ["Malta", "Malte", "Malta", "Malta"],
        "Estonia": ["Estonia", "Estonie", "Estland", "Estonia"],
        "Latvia": ["Latvia", "Lettonie", "Lettland", "Letonia"],
        "Lithuania": ["Lithuania", "Lituanie", "Litauen", "Lituania"],
        
        # Amérique
        "United States": ["United States", "USA", "États-Unis", "Vereinigte Staaten"],
        "Canada": ["Canada", "Canada", "Kanada", "Canadá"],
        "Brazil": ["Brazil", "Brésil", "Brasilien", "Brasil"],
        "Mexico": ["Mexico", "Mexique", "Mexiko", "México"],
        
        # Asie
        "Japan": ["Japan", "Japon", "Japan", "Japón"],
        "China": ["China", "Chine", "China", "China"],
        "Hong Kong": ["Hong Kong", "Hong Kong", "Hongkong", "Hong Kong"],
        "Singapore": ["Singapore", "Singapour", "Singapur", "Singapur"],
        "Australia": ["Australia", "Australie", "Australien", "Australia"],
        "South Africa": ["South Africa", "Afrique du Sud", "Südafrika", "Sudáfrica"],
    }
    
    @staticmethod
    def detect_countries_in_text(text: str) -> List[str]:
        """Détecte les pays dans un texte avec une méthode robuste"""
        detected_countries = []
        text_lower = text.lower()
        
        # Recherche directe des pays dans le texte
        for country_en, aliases in IntelligentCountryDetector.COUNTRIES_DATABASE.items():
            found = False
            
            # Vérifier chaque alias
            for alias in aliases:
                alias_lower = alias.lower()
                
                # Recherche avec contexte pour éviter les faux positifs
                # Chercher le mot entier, pas une partie d'un autre mot
                pattern = r'\b' + re.escape(alias_lower) + r'\b'
                if re.search(pattern, text_lower):
                    detected_countries.append(country_en)
                    found = True
                    break
            
            if found:
                continue
            
            # Recherche par racine commune
            country_root = country_en.lower().split()[0] if ' ' in country_en else country_en.lower()
            if len(country_root) > 3:
                pattern = r'\b' + re.escape(country_root[:4]) + r'[\w]*\b'
                matches = re.findall(pattern, text_lower)
                for match in matches:
                    if Levenshtein.ratio(country_root, match) > 0.7:
                        detected_countries.append(country_en)
                        break
        
        return list(set(detected_countries))
    
    @staticmethod
    def extract_countries_with_context(text: str) -> Dict[str, List[str]]:
        """Extrait les pays avec leur contexte d'apparition"""
        countries_with_context = {}
        text_lower = text.lower()
        
        for country_en, aliases in IntelligentCountryDetector.COUNTRIES_DATABASE.items():
            found_contexts = []
            
            for alias in aliases:
                alias_lower = alias.lower()
                pattern = r'\b' + re.escape(alias_lower) + r'\b'
                
                # Rechercher toutes les occurrences
                for match in re.finditer(pattern, text_lower):
                    start = max(0, match.start() - 50)
                    end = min(len(text), match.end() + 50)
                    context = text[start:end].strip()
                    found_contexts.append({
                        "alias": alias,
                        "context": context,
                        "position": match.start()
                    })
            
            if found_contexts:
                # Trier par position
                found_contexts.sort(key=lambda x: x["position"])
                countries_with_context[country_en] = found_contexts
        
        return countries_with_context
    
    @staticmethod
    def normalize_country_name(country: str) -> str:
        """Normalise un nom de pays vers l'anglais"""
        if not country:
            return ""
        
        country_lower = country.lower().strip()
        
        # Recherche directe dans la base de données
        for country_en, aliases in IntelligentCountryDetector.COUNTRIES_DATABASE.items():
            if country_lower == country_en.lower():
                return country_en
            
            for alias in aliases:
                if country_lower == alias.lower():
                    return country_en
        
        # Recherche par similarité
        best_match = None
        best_score = 0
        
        for country_en, aliases in IntelligentCountryDetector.COUNTRIES_DATABASE.items():
            # Comparer avec le nom anglais
            score = Levenshtein.ratio(country_lower, country_en.lower())
            if score > best_score and score > 0.7:
                best_score = score
                best_match = country_en
            
            # Comparer avec les alias
            for alias in aliases:
                score = Levenshtein.ratio(country_lower, alias.lower())
                if score > best_score and score > 0.7:
                    best_score = score
                    best_match = country_en
        
        return best_match if best_match else country.title()

# ==================== FONCTIONS D'AIDE POUR LA DÉTECTION DE L'ISIN ====================
def extract_isin_from_text(text: str) -> List[str]:
    """
    Extrait TOUS les codes ISIN valides du texte.
    Format ISIN: 2 lettres suivies de 10 caractères alphanumériques
    """
    isins = []
    
    # Pattern pour détecter ISIN (2 lettres + 9 chiffres ou lettres + 1 caractère de contrôle)
    isin_pattern = r'\b[A-Z]{2}[A-Z0-9]{9}[0-9]\b'
    matches = re.findall(isin_pattern, text)
    
    for match in matches:
        if TextUtils.is_valid_isin(match):
            isins.append(match)
    
    # Recherche de motifs moins stricts
    isin_keyword_pattern = r'(?:ISIN|isin)[\s:]*([A-Z0-9]{12})'
    matches_keyword = re.findall(isin_keyword_pattern, text, re.IGNORECASE)
    
    for match in matches_keyword:
        if TextUtils.is_valid_isin(match):
            isin = match.upper()
            if isin not in isins:
                isins.append(isin)
    
    # Recherche dans les codes de 12 caractères alphanumériques avec validation
    alnum_12_pattern = r'\b([A-Z0-9]{12})\b'
    matches_alnum = re.findall(alnum_12_pattern, text)
    
    for candidate in matches_alnum:
        if TextUtils.is_valid_isin(candidate) and candidate not in isins:
            isins.append(candidate)
    
    return isins

def search_isin_in_slides(slides_data: List[Dict]) -> List[str]:
    """
    Recherche TOUS les ISIN valides dans tous les slides
    """
    all_isins = []
    
    for slide in slides_data:
        slide_isins = extract_isin_from_text(slide["text"])
        for isin in slide_isins:
            if isin not in all_isins and TextUtils.is_valid_isin(isin):
                all_isins.append(isin)
    
    # Recherche dans les titres également
    for slide in slides_data:
        for title in slide.get("titles", []):
            title_isins = extract_isin_from_text(title)
            for isin in title_isins:
                if isin not in all_isins and TextUtils.is_valid_isin(isin):
                    all_isins.append(isin)
    
    return all_isins

# ==================== ÉTAT COMBINÉ ====================
class CombinedAgentState(TypedDict):
    # Entrée utilisateur
    pptx_path: str
    glossaires_json_path: str
    
    # Contenu du document
    presentation: Any
    slides_data: List[Dict]
    slides_content: List[Dict]
    full_text: str
    error: Optional[str]
    
    # Contexte identifié
    detected_fund: Dict[str, Any]
    detected_lang: str
    detected_target: str
    detected_countries: List[str]
    normalized_countries: List[str]
    
    # Ressources chargées
    ref_fund_data: Dict[str, Any]
    all_registered_countries: Dict[str, Any]
    required_disclaimers: List[Dict]
    required_glossaires: List[Dict]
    glossaires_attendus: List[str]
    
    # Résultats d'analyse
    registration_check: Dict[str, Any]
    disclaimer_check: Dict[str, Any]
    glossary_check: Dict[str, Any]
    source_check: Dict[str, Any]
    characteristics_check: Dict[str, Any]
    
    # Métriques et scores
    resultats: Dict[str, Any]
    metriques: Dict[str, Any]
    scores_alignement: List[Dict]
    qualite_reponses: List[Dict]
    
    # Index et positions
    current_slide_idx: int
    annexe_slide_idx: int
    characteristics_slide_idx: int
    
    # Rapports
    report: Dict[str, Any]
    erreurs: List[str]
    
    # Métriques LLM
    llm_call_count: int
    llm_warning: Optional[str]

# ==================== COLLECTEUR DE MÉTRIQUES AMÉLIORÉ ====================
class EnhancedMetricsCollector:
    def __init__(self):
        self.start_time = None
        self.end_time = None
        self.slides_traites = 0
        self.glossaires_detectes = 0
        self.disclaimers_detectes = 0
        self.sources_verifiees = 0
        self.llm_calls = 0
        self.erreurs = 0
        self.temps_llm_total = 0
        self.temps_par_slide = []
        self.scores_similarite = []
        self.scores_rouge = []
        self.scores_readability = []
        self.coverage_scores = []
        self.compliance_scores = []
        self.tokens_utilises = 0
        
    def start(self):
        self.start_time = time.time()
    
    def stop(self):
        self.end_time = time.time()
    
    def add_slide_time(self, time_sec: float):
        self.temps_par_slide.append(time_sec)
    
    def add_coverage_score(self, score: float):
        self.coverage_scores.append(score)
    
    def add_compliance_score(self, score: float):
        self.compliance_scores.append(score)
    
    def add_similarity_score(self, score: float):
        self.scores_similarite.append(score)
    
    def add_rouge_score(self, score: float):
        self.scores_rouge.append(score)
    
    def add_readability_score(self, score: float):
        self.scores_readability.append(score)
    
    def add_llm_call(self, duration: float, tokens_estimes: int = 0):
        """Ajoute les métriques d'un appel LLM"""
        self.llm_calls += 1
        self.temps_llm_total += duration
        self.tokens_utilises += tokens_estimes
    
    def to_dict(self):
        avg_similarity = np.mean(self.scores_similarite) if self.scores_similarite else 0
        avg_rouge = np.mean(self.scores_rouge) if self.scores_rouge else 0
        avg_readability = np.mean(self.scores_readability) if self.scores_readability else 0
        avg_coverage = np.mean(self.coverage_scores) if self.coverage_scores else 0
        avg_compliance = np.mean(self.compliance_scores) if self.compliance_scores else 0
        
        slide_times = self.temps_par_slide
        avg_slide_time = np.mean(slide_times) if slide_times else 0
        
        return {
            "duree_execution_secondes": round(self.end_time - self.start_time, 2) if self.end_time else None,
            "temps_llm_total": round(self.temps_llm_total, 2),
            "temps_llm_moyen": round(self.temps_llm_total / max(self.llm_calls, 1), 2),
            "temps_slide_moyen": round(avg_slide_time, 2),
            "slides_traites": self.slides_traites,
            "glossaires_detectes": self.glossaires_detectes,
            "disclaimers_detectes": self.disclaimers_detectes,
            "sources_verifiees": self.sources_verifiees,
            "appels_llm": self.llm_calls,
            "tokens_estimes": self.tokens_utilises,
            "erreurs": self.erreurs,
            "scores_couverture": {
                "moyenne": round(avg_coverage, 3),
                "min": round(min(self.coverage_scores), 3) if self.coverage_scores else 0,
                "max": round(max(self.coverage_scores), 3) if self.coverage_scores else 0
            },
            "scores_conformite": {
                "moyenne": round(avg_compliance, 3),
                "min": round(min(self.compliance_scores), 3) if self.compliance_scores else 0,
                "max": round(max(self.compliance_scores), 3) if self.compliance_scores else 0
            },
            "metriques_qualite": {
                "similarite_cosine_moyenne": round(avg_similarity, 3),
                "score_rouge_moyen": round(avg_rouge, 3),
                "lisibilite_moyenne": round(avg_readability, 2),
                "distribution_similarite": {
                    "min": round(min(self.scores_similarite), 3) if self.scores_similarite else 0,
                    "max": round(max(self.scores_similarite), 3) if self.scores_similarite else 0,
                    "median": round(np.median(self.scores_similarite), 3) if self.scores_similarite else 0,
                    "ecart_type": round(np.std(self.scores_similarite), 3) if self.scores_similarite else 0
                }
            },
            "efficacite": {
                "slides_par_seconde": round(self.slides_traites / (self.end_time - self.start_time), 3) if self.end_time and self.start_time else 0,
                "taux_succes": 1 - (self.erreurs / max(self.slides_traites, 1)),
                "densite_information": round(self.tokens_utilises / max(self.slides_traites, 1), 1)
            }
        }

# ==================== UTILITAIRES DE SIMILARITÉ ====================
class SimilarityCalculator:
    @staticmethod
    def cosine_similarity(text1: str, text2: str) -> float:
        """Calcule la similarité cosinus entre deux textes"""
        if not text1 or not text2:
            return 0.0
        
        vectorizer = TfidfVectorizer().fit_transform([text1, text2])
        vectors = vectorizer.toarray()
        
        cos_sim = cosine_similarity([vectors[0]], [vectors[1]])[0][0]
        return float(cos_sim)

# ==================== UTILITAIRES GÉNÉRAUX ====================
def extract_text_from_shape(shape):
    """Extrait le texte d'une forme PPTX"""
    try:
        if hasattr(shape, "text"):
            return shape.text.strip()
        if hasattr(shape, "text_frame"):
            return shape.text_frame.text.strip()
        if hasattr(shape, "has_table") and shape.has_table:
            text = ""
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text_frame:
                        text += cell.text_frame.text + " "
            return text.strip()
    except:
        pass
    return ""

def is_text_bold(run):
    """Vérifie si le texte est en gras"""
    try:
        return run.font.bold == True
    except:
        return False

def extract_footer_content(slide):
    """Extrait le contenu du footer (disclaimer et source)"""
    footer_data = {
        "disclaimer": {"texte": "", "en_gras": False, "langue": ""},
        "source": {"texte": "", "vide": True, "contient_seulement_date": False}
    }
    
    try:
        slide_height = 540  # Valeur par défaut pour PowerPoint
        
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                
                # Vérifier si la forme est en bas du slide (position > 70% de la hauteur)
                if shape.top > slide_height * 0.7:
                    # Recherche de texte en gras pour le disclaimer
                    has_bold = False
                    disclaimer_text = ""
                    
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if is_text_bold(run):
                                has_bold = True
                                disclaimer_text += run.text
                            elif has_bold and disclaimer_text:
                                disclaimer_text += run.text
                    
                    if has_bold and disclaimer_text:
                        footer_data["disclaimer"]["texte"] = disclaimer_text
                        footer_data["disclaimer"]["en_gras"] = True
                    
                    # Chercher la source
                    source_match = re.search(r'(?:source|SOURCE)\s*:?\s*(.+)', text, re.IGNORECASE)
                    if source_match:
                        source_text = source_match.group(1).strip()
                        footer_data["source"]["texte"] = source_text
                        footer_data["source"]["vide"] = len(source_text) == 0
                    elif not footer_data["disclaimer"]["texte"]:
                        if len(text) < 200:
                            footer_data["source"]["texte"] = text
                            footer_data["source"]["vide"] = False
        
        # Si aucun footer trouvé en bas, chercher dans tout le slide pour le disclaimer
        if not footer_data["disclaimer"]["texte"]:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    text = shape.text_frame.text.strip()
                    if not text or len(text) < 10:
                        continue
                    
                    has_bold = False
                    disclaimer_text = ""
                    
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if is_text_bold(run):
                                has_bold = True
                                disclaimer_text += run.text
                    
                    if has_bold and disclaimer_text and len(disclaimer_text) > 20:
                        footer_data["disclaimer"]["texte"] = disclaimer_text
                        footer_data["disclaimer"]["en_gras"] = True
                        break
                        
    except Exception as e:
        print(f"Erreur extraction footer: {e}")
    
    return footer_data

# ==================== AGENT OPTIMISÉ ====================
class OptimizedDisclaimerGlossaryAgent:
    def __init__(self, max_llm_calls: int = 50):
        """
        Agent optimisé avec gestion intelligente des appels LLM
        
        Args:
            max_llm_calls: Nombre maximum d'appels LLM avant avertissement (non limitatif)
        """
        self.metrics = EnhancedMetricsCollector()
        self.metrics.start()
        self.llm_client = LlamaClient()
        self.max_llm_calls = max_llm_calls
        self.llm_call_count = 0
        self.llm_warning_issued = False
        self.text_utils = TextUtils()
        self.country_detector = IntelligentCountryDetector()
        
    def _call_llm(self, prompt: str, system_message: str = None, temperature: float = 0.3, 
                  max_tokens: int = 1000, json_mode: bool = False) -> Tuple[Any, Dict]:
        """Appelle le LLM via le client centralisé avec comptage"""
        start_time = time.time()
        self.llm_call_count += 1
        
        # Avertissement si on dépasse le seuil
        if self.llm_call_count > self.max_llm_calls and not self.llm_warning_issued:
            print(f"⚠️  ATTENTION: {self.llm_call_count} appels LLM effectués (seuil: {self.max_llm_calls})")
            print("   L'analyse continue sans limitation des performances")
            self.llm_warning_issued = True
        
        print(f"📞 Appel LLM #{self.llm_call_count}")
        
        try:
            if json_mode:
                system_prompt = system_message or """Tu es un expert en analyse de documents financiers. 
                Réponds UNIQUEMENT en JSON valide. Ne rajoute aucun texte avant ou après le JSON."""
                result = self.llm_client.generate_json_response(system_prompt, prompt, temperature=temperature)
                
                # Estimation des tokens
                if isinstance(result, dict):
                    json_str = json.dumps(result)
                    tokens_estimes = len(json_str.split()) * 1.3
                else:
                    tokens_estimes = 0
                
                response_time = time.time() - start_time
                self.metrics.add_llm_call(response_time, tokens_estimes)
                
                return result, {
                    "temps_reponse": round(response_time, 2),
                    "mode": "json",
                    "call_number": self.llm_call_count
                }
            else:
                system_message = system_message or """Tu es un assistant expert en analyse de documents financiers. 
                Tes réponses doivent être :
                1. Concises et précises
                2. Structurées avec des marqueurs clairs
                3. Basées uniquement sur les informations fournies
                4. En français professionnel sauf indication contraire"""
                
                messages = [
                    {"role": "system", "content": system_message},
                    {"role": "user", "content": prompt}
                ]
                
                response = self.llm_client.generate_response(
                    messages, 
                    temperature=temperature, 
                    max_tokens=max_tokens
                )
                
                response_time = time.time() - start_time
                tokens_estimes = len(response.split()) * 1.3 if response else 0
                self.metrics.add_llm_call(response_time, tokens_estimes)
                
                return response, {
                    "temps_reponse": round(response_time, 2),
                    "longueur_reponse": len(response) if response else 0,
                    "mode": "text",
                    "call_number": self.llm_call_count
                }
                
        except Exception as e:
            response_time = time.time() - start_time
            self.metrics.erreurs += 1
            self.metrics.add_llm_call(response_time, 0)
            
            error_msg = f"ERREUR: {str(e)}"
            if json_mode:
                return {"error": error_msg}, {"erreur": str(e), "call_number": self.llm_call_count}
            else:
                return error_msg, {"erreur": str(e), "call_number": self.llm_call_count}
    
    def _load_resources(self) -> Dict[str, Any]:
        """Charge les ressources JSON en mémoire"""
        resources = {}
        try:
            if os.path.exists(FOND_REGISTRED_PATH):
                with open(FOND_REGISTRED_PATH, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                    resources['fonds'] = data.get('fonds', [])
                    print(f"✅ {len(resources['fonds'])} fonds chargés")
            else:
                resources['fonds'] = []
                print(f"⚠️  Fichier non trouvé: {FOND_REGISTRED_PATH}")
                
            if os.path.exists(DISCLAIMERS_PATH):
                with open(DISCLAIMERS_PATH, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                    resources['disclaimers'] = data.get('disclaimers', [])
                    print(f"✅ {len(resources['disclaimers'])} disclaimers chargés")
            else:
                resources['disclaimers'] = []
                print(f"⚠️  Fichier non trouvé: {DISCLAIMERS_PATH}")
                
            if os.path.exists(GLOSSAIRES_PATH):
                with open(GLOSSAIRES_PATH, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                    resources['glossaires'] = data
                    print(f"✅ Glossaires chargés")
            else:
                resources['glossaires'] = {}
                print(f"⚠️  Fichier non trouvé: {GLOSSAIRES_PATH}")
                
        except Exception as e:
            print(f"❌ Erreur chargement ressources: {e}")
            resources = {'fonds': [], 'disclaimers': [], 'glossaires': {}}
        
        return resources

    # ==================== NŒUDS DU GRAPHE OPTIMISÉS ====================
    
    def parse_document_node(self, state: CombinedAgentState) -> CombinedAgentState:
        """Parse le document PPTX de manière intelligente"""
        print(f"📄 Parsing intelligent du document: {os.path.basename(state['pptx_path'])}")
        
        try:
            state["presentation"] = Presentation(state["pptx_path"])
            state["slides_data"] = []
            state["slides_content"] = []
            state["full_text"] = ""
            state["erreurs"] = []
            state["llm_call_count"] = 0
            state["llm_warning"] = None
            
            for i, slide in enumerate(state["presentation"].slides):
                slide_start_time = time.time()
                
                slide_text = ""
                slide_titles = []
                
                # Extraction intelligente du texte
                for shape in slide.shapes:
                    text = extract_text_from_shape(shape)
                    if text:
                        # Identifier les titres (généralement en haut et en gras)
                        try:
                            if hasattr(shape, "text_frame"):
                                for paragraph in shape.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        if run.font.bold and len(run.text.strip()) > 3:
                                            slide_titles.append(run.text.strip())
                        except:
                            pass
                        
                        slide_text += text + "\n"
                
                # Nettoyage et structuration
                slide_text_clean = re.sub(r'\s+', ' ', slide_text).strip()
                
                slide_data = {
                    "index": i,
                    "text": slide_text_clean,
                    "titles": list(set(slide_titles)),
                    "has_characteristics": False,
                    "has_glossary": False,
                    "footer": extract_footer_content(slide)
                }
                
                state["slides_data"].append(slide_data)
                state["slides_content"].append({"index": i, "text": slide_text_clean})
                state["full_text"] += f"--- SLIDE {i+1} ---\n{slide_text_clean}\n\n"
                
                slide_time = time.time() - slide_start_time
                self.metrics.add_slide_time(slide_time)
            
            self.metrics.slides_traites = len(state["slides_data"])
            print(f"  ✅ Parsing réussi: {len(state['slides_data'])} slides analysés")
            
        except Exception as e:
            print(f"❌ Erreur parsing: {e}")
            state["error"] = str(e)
            self.metrics.erreurs += 1
            
        return state

    def load_glossaires_node(self, state: CombinedAgentState) -> CombinedAgentState:
        """Charge les glossaires depuis le JSON spécifié"""
        try:
            with open(state["glossaires_json_path"], 'r', encoding='utf-8') as f:
                data = json.load(f)
                
                glossaires_list = []
                if isinstance(data, dict):
                    # Format: {"glossaire_1": {"nom": "...", "description": "..."}}
                    for key, value in data.items():
                        if isinstance(value, dict) and "nom" in value:
                            glossaires_list.append(value["nom"])
                        elif isinstance(value, str):
                            glossaires_list.append(value)
                elif isinstance(data, list):
                    # Format: ["terme1", "terme2", ...]
                    glossaires_list = data
                
                state["glossaires_attendus"] = [g.strip() for g in glossaires_list if g.strip()]
            
            print(f"✓ {len(state['glossaires_attendus'])} glossaires chargés")
            
        except Exception as e:
            print(f"❌ Erreur chargement glossaires: {e}")
            state["erreurs"].append(f"Erreur chargement glossaires: {e}")
            self.metrics.erreurs += 1
        
        return state

    def identify_fund_and_context_node(self, state: CombinedAgentState) -> CombinedAgentState:
        """Identification intelligente du fonds et du contexte via LLM"""
        print("🔍 Identification intelligente du fonds et contexte...")
        
        # RECHERCHE PRÉLIMINAIRE DE TOUS LES ISINS AVEC REGEX
        all_isins = search_isin_in_slides(state["slides_data"])
        print(f"  🔍 Recherche regex ISIN: {len(all_isins)} ISIN(s) valide(s) - {all_isins}")
        
        # Détection des pays avec la méthode améliorée
        detected_countries = self.country_detector.detect_countries_in_text(state["full_text"])
        print(f"  🌍 Détection de pays: {len(detected_countries)} pays détectés - {detected_countries}")
        
        # Normaliser les pays détectés
        normalized_countries = []
        for country in detected_countries:
            normalized = self.country_detector.normalize_country_name(country)
            if normalized and normalized not in normalized_countries:
                normalized_countries.append(normalized)
        
        # Utiliser le LLM pour extraire les informations manquantes
        system_prompt = """Tu es un expert en analyse de documents financiers. 
        Analyse le document pour extraire les informations suivantes.
        Réponds UNIQUEMENT en JSON avec cette structure exacte:
        {
            "nom_fonds": "Nom exact du fonds",
            "classe_part": "Classe de parts si mentionnée",
            "langue": "fr/en/de/es/it",
            "cible": "Retail/Professionnel/Both",
            "caracteristiques_principales": ["caractéristique1", "caractéristique2"],
            "slide_caracteristiques": numéro du slide contenant les caractéristiques (0-index),
            "note_detection": "notes sur la détection"
        }"""
        
        user_prompt = f"""Analyse ce document PowerPoint pour extraire les informations financières.
        
Document complet (extrait):
{state['full_text'][:5000]}... [document tronqué pour économiser des tokens]

Instructions:
1. Identifie le nom exact du fonds
2. Identifie la classe de parts si mentionnée
3. Détermine la langue principale du document
4. Identifie la cible (Retail, Professionnel, ou Both)
5. Identifie les caractéristiques principales du fonds
6. Trouve le numéro du slide contenant les caractéristiques (0-index)"""

        try:
            response, metadata = self._call_llm(user_prompt, system_prompt, temperature=0.1, json_mode=True)
            state["llm_call_count"] = self.llm_call_count
            
            if isinstance(response, dict) and "error" not in response:
                # Extraire les informations
                fund_name = response.get("nom_fonds", "Inconnu")
                classe_part = response.get("classe_part", None)
                language = response.get("langue", "fr").lower()
                target = response.get("cible", "Professionnel")
                characteristics_slide = response.get("slide_caracteristiques", -1)
                
                state["detected_fund"] = {
                    "name": fund_name,
                    "isins": all_isins,
                    "classe_part": classe_part
                }
                state["detected_lang"] = language
                state["detected_target"] = target
                state["detected_countries"] = detected_countries
                state["normalized_countries"] = normalized_countries
                state["characteristics_slide_idx"] = characteristics_slide
                
                print(f"  ✅ FONDS IDENTIFIÉ: {fund_name}")
                print(f"     ISINs valides: {all_isins if all_isins else 'Non trouvé'}")
                print(f"     Langue: {language}, Cible: {target}")
                print(f"     Pays détectés: {detected_countries}")
                print(f"     Pays normalisés: {normalized_countries}")
                print(f"     Slide caractéristiques: {characteristics_slide if characteristics_slide != -1 else 'Non trouvé'}")
                
                # Chargement des ressources
                resources = self._load_resources()
                
                # Recherche de TOUS les fonds correspondant aux ISINs détectés
                found_funds = []
                all_registered_countries = {}
                
                for isin in all_isins:
                    for fund in resources['fonds']:
                        if fund.get('isin') == isin and fund not in found_funds:
                            found_funds.append(fund)
                            # Collecter les pays d'enregistrement de tous les fonds trouvés
                            pays_enregistrement = fund.get("pays_enregistrement", {})
                            for pays, info in pays_enregistrement.items():
                                # Normaliser le nom du pays
                                normalized_pays = self.country_detector.normalize_country_name(pays)
                                if normalized_pays:
                                    all_registered_countries[normalized_pays] = {
                                        "original": pays,
                                        "info": info,
                                        "normalized": normalized_pays
                                    }
                
                # Si aucun fonds trouvé par ISIN, recherche par nom
                if not found_funds:
                    search_name = fund_name.lower()
                    for fund in resources['fonds']:
                        fund_name_lower = fund.get('nom_fonds', '').lower()
                        # Recherche approximative
                        if (search_name in fund_name_lower or 
                            fund_name_lower in search_name or
                            Levenshtein.ratio(search_name, fund_name_lower) > 0.8):
                            found_funds.append(fund)
                            # Collecter les pays d'enregistrement
                            pays_enregistrement = fund.get("pays_enregistrement", {})
                            for pays, info in pays_enregistrement.items():
                                normalized_pays = self.country_detector.normalize_country_name(pays)
                                if normalized_pays:
                                    all_registered_countries[normalized_pays] = {
                                        "original": pays,
                                        "info": info,
                                        "normalized": normalized_pays
                                    }
                
                # Stocker le premier fonds trouvé (ou un dict vide) et tous les pays enregistrés
                state["ref_fund_data"] = found_funds[0] if found_funds else {}
                state["all_registered_countries"] = all_registered_countries
                
                # Filtrer les disclaimers requis
                state["required_disclaimers"] = [
                    d for d in resources['disclaimers']
                    if d.get('langue') == language and d.get('cible') == target
                ]
                
                # Si aucun disclaimer trouvé pour cette combinaison, utiliser ceux de la langue
                if not state["required_disclaimers"]:
                    state["required_disclaimers"] = [
                        d for d in resources['disclaimers']
                        if d.get('langue') == language
                    ]
                
                # Glossaires obligatoires
                state["required_glossaires"] = resources['glossaires']
                
                if found_funds:
                    print(f"  ✅ {len(found_funds)} fond(s) trouvé(s) dans la base")
                    print(f"     Pays enregistrés: {list(all_registered_countries.keys())}")
                else:
                    print(f"  ⚠️  Fonds non trouvé dans la base")
                    
            else:
                print("❌ Erreur dans l'analyse LLM, utilisation des valeurs par défaut")
                state["detected_fund"] = {
                    "name": "Inconnu", 
                    "isins": all_isins, 
                    "classe_part": None
                }
                state["detected_lang"] = "fr"
                state["detected_target"] = "Professionnel"
                state["detected_countries"] = detected_countries
                state["normalized_countries"] = normalized_countries
                state["characteristics_slide_idx"] = -1
                state["ref_fund_data"] = {}
                state["all_registered_countries"] = {}
                state["required_disclaimers"] = []
                state["required_glossaires"] = {}
                
        except Exception as e:
            print(f"❌ Erreur identification contexte: {e}")
            state["detected_fund"] = {
                "name": "Inconnu", 
                "isins": all_isins, 
                "classe_part": None
            }
            state["detected_lang"] = "fr"
            state["detected_target"] = "Professionnel"
            state["detected_countries"] = detected_countries
            state["normalized_countries"] = normalized_countries
            state["characteristics_slide_idx"] = -1
            state["ref_fund_data"] = {}
            state["all_registered_countries"] = {}
            state["required_disclaimers"] = []
            state["required_glossaires"] = {}
        
        return state

    def verify_disclaimers_intelligently_node(self, state: CombinedAgentState) -> CombinedAgentState:
        """Vérification intelligente des disclaimers et de leur alignement - OPTIMISÉ"""
        print("⚖️ Vérification intelligente des disclaimers...")
        
        resultats = {
            "obligatoires": {"conformes": [], "non_conformes": [], "absents": []},
            "alignement": {"conformes": [], "non_conformes": [], "absents": []},
            "scores": [],
            "note": "Les disclaimers obligatoires sont spécifiques à la langue et à la cible du document."
        }
        
        required_disclaimers = state.get("required_disclaimers", [])
        language = state.get("detected_lang", "fr")
        target = state.get("detected_target", "Professionnel")
        
        print(f"  📝 Recherche de {len(required_disclaimers)} disclaimers pour {language}/{target}")
        
        if not required_disclaimers:
            print("  ⚠️  Aucun disclaimer requis trouvé pour cette combinaison langue/cible")
            # Charger tous les disclaimers pour analyse
            resources = self._load_resources()
            required_disclaimers = resources.get('disclaimers', [])
            print(f"  🔄 Utilisation de {len(required_disclaimers)} disclaimers génériques pour l'analyse")
        
        # STRATÉGIE OPTIMISÉE: Analyser tous les disclaimers obligatoires en UN SEUL appel LLM
        if required_disclaimers:
            system_prompt = f"""Tu es un expert en vérification de documents financiers en {language}.
            Vérifie quels disclaimers obligatoires sont présents dans le document.
            Réponds UNIQUEMENT en JSON:
            {{
                "disclaimers_verifies": [
                    {{
                        "titre": "titre du disclaimer",
                        "present": true/false,
                        "similarity_score": 0.0-1.0,
                        "location": "slide number or 'document'"
                    }}
                ],
                "summary": {{
                    "total_required": nombre,
                    "found_count": nombre,
                    "average_similarity": moyenne
                }},
                "note": "notes sur la vérification"
            }}"""
            
            user_prompt = f"""Disclaimers obligatoires requis ({len(required_disclaimers)}):

{json.dumps(required_disclaimers, ensure_ascii=False, indent=2)}

Document complet (extrait):
{state['full_text'][:3000]}... [document tronqué]

Pour chaque disclaimer obligatoire:
1. Vérifie s'il est présent (identique ou similaire)
2. Donne un score de similarité (0-1)
3. Indique la localisation (numéro de slide ou 'document')"""

            try:
                response, _ = self._call_llm(user_prompt, system_prompt, temperature=0.1, json_mode=True)
                
                if isinstance(response, dict):
                    disclaimers_verifies = response.get("disclaimers_verifies", [])
                    
                    for dv in disclaimers_verifies:
                        titre = dv.get("titre", "")
                        present = dv.get("present", False)
                        similarity = dv.get("similarity_score", 0)
                        location = dv.get("location", "inconnu")
                        
                        if present and similarity >= 0.7:
                            resultats["obligatoires"]["conformes"].append({
                                "titre": titre,
                                "similarity": similarity,
                                "location": location
                            })
                        elif present:
                            resultats["obligatoires"]["non_conformes"].append({
                                "titre": titre,
                                "similarity": similarity,
                                "raison": f"Similarité trop faible ({similarity:.2f})"
                            })
                        else:
                            resultats["obligatoires"]["absents"].append({
                                "titre": titre,
                                "raison": "Non détecté dans le document"
                            })
                            
                    # Mettre à jour la note
                    if "note" in response:
                        resultats["note"] = response["note"]
                        
            except Exception as e:
                print(f"⚠️ Erreur vérification disclaimers groupés: {e}")
                # Fallback: vérification simple par recherche textuelle
                full_text_lower = state.get("full_text", "").lower()
                for req in required_disclaimers:
                    content = req.get('contenu', '').lower()
                    title = req.get('titre', '')
                    
                    if content and content in full_text_lower:
                        resultats["obligatoires"]["conformes"].append({
                            "titre": title,
                            "similarity": 0.9,
                            "location": "Document"
                        })
                    else:
                        resultats["obligatoires"]["absents"].append({
                            "titre": title,
                            "raison": "Non trouvé par recherche textuelle"
                        })
        
        # STRATÉGIE OPTIMISÉE: Analyser l'alignement des disclaimers par groupe de slides
        slides_with_disclaimers = []
        for idx, slide_data in enumerate(state["slides_data"]):
            disclaimer_texte = slide_data["footer"]["disclaimer"]["texte"]
            
            if disclaimer_texte and len(disclaimer_texte) > 10:
                slides_with_disclaimers.append({
                    "index": idx,
                    "disclaimer": disclaimer_texte,
                    "slide_content": slide_data["text"][:500]  # Limiter la taille
                })
                self.metrics.disclaimers_detectes += 1
        
        print(f"  📊 {len(slides_with_disclaimers)} slides avec des disclaimers détectés")
        
        # Analyser l'alignement par groupes de 5 slides maximum
        if slides_with_disclaimers:
            for i in range(0, len(slides_with_disclaimers), 5):
                group = slides_with_disclaimers[i:i+5]
                
                system_prompt = """Évalue l'alignement entre le contenu des slides et leurs disclaimers.
                Réponds UNIQUEMENT en JSON:
                {
                    "alignments": [
                        {
                            "slide_index": numéro,
                            "alignment_score": 0.0-1.0,
                            "relevant": true/false,
                            "justification": "brief explanation"
                        }
                    ],
                    "summary": {
                        "total_evaluated": nombre,
                        "average_alignment": moyenne
                    }
                }"""
                
                user_prompt = f"""Évalue l'alignement pour les slides suivants:

{json.dumps(group, ensure_ascii=False, indent=2)}

Pour chaque slide:
1. Score d'alignement entre le contenu et le disclaimer (0-1)
2. Pertinence (true/false) - le disclaimer est-il approprié pour ce contenu?
3. Justification brève"""
                
                try:
                    response, _ = self._call_llm(user_prompt, system_prompt, temperature=0.2, json_mode=True)
                    
                    if isinstance(response, dict):
                        alignments = response.get("alignments", [])
                        
                        for alignment in alignments:
                            idx = alignment.get("slide_index")
                            alignment_score = alignment.get("alignment_score", 0)
                            relevant = alignment.get("relevant", False)
                            
                            self.metrics.add_similarity_score(alignment_score)
                            
                            resultats["scores"].append({
                                "slide": idx + 1,
                                "alignment_score": round(alignment_score, 3),
                                "relevant": relevant,
                                "justification": alignment.get("justification", "")
                            })
                            
                            if alignment_score >= 0.5:
                                resultats["alignement"]["conformes"].append({
                                    "slide": idx + 1,
                                    "score": round(alignment_score, 3),
                                    "justification": alignment.get("justification", "")
                                })
                            else:
                                resultats["alignement"]["non_conformes"].append({
                                    "slide": idx + 1,
                                    "score": round(alignment_score, 3),
                                    "raison": f"Alignement insuffisant: {alignment.get('justification', '')}"
                                })
                except Exception as e:
                    print(f"⚠️ Erreur alignement groupe slides: {e}")
                    # Marquer comme absent en cas d'erreur
                    for slide in group:
                        resultats["alignement"]["absents"].append({
                            "slide": slide["index"] + 1,
                            "raison": "Erreur d'analyse d'alignement"
                        })
        
        # Marquer les slides sans disclaimers comme absents
        for idx, slide_data in enumerate(state["slides_data"]):
            disclaimer_texte = slide_data["footer"]["disclaimer"]["texte"]
            if not disclaimer_texte or len(disclaimer_texte) <= 10:
                resultats["alignement"]["absents"].append({
                    "slide": idx + 1,
                    "raison": "Pas de disclaimer détecté"
                })
        
        # Calcul des scores
        obligatoire_score = len(resultats["obligatoires"]["conformes"]) / max(len(required_disclaimers), 1)
        alignement_score = len(resultats["alignement"]["conformes"]) / max(len(state["slides_data"]), 1)
        
        state["disclaimer_check"] = {
            **resultats,
            "obligatoire_score": obligatoire_score,
            "alignement_score": alignement_score,
            "global_score": (obligatoire_score + alignement_score) / 2,
            "detected_language": state.get("detected_lang"),
            "detected_target": state.get("detected_target"),
            "total_slides": len(state["slides_data"]),
            "slides_with_disclaimers": len(slides_with_disclaimers)
        }
        
        self.metrics.add_coverage_score(obligatoire_score)
        self.metrics.add_compliance_score(alignement_score)
        
        return state

    def verify_glossaries_intelligently_node(self, state: CombinedAgentState) -> CombinedAgentState:
        """Vérification intelligente des glossaires - CORRIGÉE"""
        print("📖 Vérification intelligente des glossaires...")
        
        resultats = {
            "obligatoires": {"conformes": [], "non_conformes": [], "absents": []},
            "utilisation": {"conformes": [], "non_conformes": [], "absents": []},
            "note": "Les 23 premiers termes sont les glossaires obligatoires. Les autres sont spécifiques au document."
        }
        
        # Glossaires obligatoires (23)
        required_glossaires = state.get("required_glossaires", {})
        
        # Trouver le slide glossaire intelligemment
        glossary_slide_idx = -1
        glossary_text = ""
        
        for idx, slide_data in enumerate(state["slides_data"]):
            slide_text_lower = slide_data["text"].lower()
            # Recherche par mots-clés étendus
            glossary_keywords = ["glossaire", "glossary", "définitions", "definitions", 
                               "vocabulaire", "vocabulary", "annexe", "appendix", "lexique"]
            if any(keyword in slide_text_lower for keyword in glossary_keywords):
                glossary_slide_idx = idx
                glossary_text = slide_data["text"]
                state["slides_data"][idx]["has_glossary"] = True
                break
        
        if glossary_slide_idx == -1:
            # Si pas trouvé, utiliser les 3 derniers slides
            last_slides = state["slides_data"][-3:] if len(state["slides_data"]) > 3 else state["slides_data"]
            for slide_data in last_slides:
                glossary_text += slide_data["text"] + "\n"
        
        # Extraire les termes obligatoires du cache des glossaires
        all_required_terms = []
        
        if isinstance(required_glossaires, dict):
            # Essayer différents formats de données
            if "glossaires" in required_glossaires:
                # Format: {"glossaires": [{"nom": "..."}, ...]}
                for item in required_glossaires.get("glossaires", []):
                    if isinstance(item, dict) and "nom" in item:
                        term = item.get("nom", "").strip()
                        if term and term not in all_required_terms:
                            all_required_terms.append(term)
            else:
                # Parcourir toutes les clés et valeurs
                for key, value in required_glossaires.items():
                    term = None
                    if isinstance(value, dict) and "nom" in value:
                        term = value.get("nom", "").strip()
                    elif isinstance(value, str):
                        term = value.strip()
                    else:
                        term = str(key).strip()
                    
                    if term and term not in all_required_terms:
                        all_required_terms.append(term)
        
        # Limiter aux 23 premiers termes obligatoires
        required_terms_23 = all_required_terms[:23]
        
        # Glossaires attendus du fichier spécifié
        glossaires_attendus = state.get("glossaires_attendus", [])
        
        print(f"  📚 {len(required_terms_23)} termes obligatoires à vérifier")
        print(f"  📖 {len(glossaires_attendus)} glossaires spécifiques attendus")
        
        # STRATÉGIE OPTIMISÉE: Vérifier les glossaires obligatoires en un seul appel
        if required_terms_23 and glossary_text:
            system_prompt = """Tu es un expert en analyse de glossaires financiers.
            Identifie les termes du glossaire présents dans le texte.
            Réponds UNIQUEMENT en JSON:
            {
                "terms_found": [
                    {
                        "term": "terme exact",
                        "found": true/false,
                        "confidence": 0.0-1.0,
                        "location": "où trouvé dans le texte"
                    }
                ],
                "summary": {
                    "total_searched": nombre,
                    "found_count": nombre,
                    "coverage_rate": pourcentage
                }
            }"""
            
            user_prompt = f"""Dans le texte suivant, identifie quels termes de glossaire financier sont présents:

TEXTE GLOSSAIRE:
{glossary_text[:2000]}... [texte tronqué]

LISTE DES 23 TERMES OBLIGATOIRES À RECHERCHER:
{json.dumps(required_terms_23, ensure_ascii=False, indent=2)}

Instructions:
1. Pour chaque terme, indique s'il est présent (found: true/false)
2. Donne un niveau de confiance (confidence: 0-1)
3. Indique où il a été trouvé dans le texte
4. Sois large dans la recherche (synonymes, variantes, formes plurielles)"""

            try:
                response, _ = self._call_llm(user_prompt, system_prompt, temperature=0.1, json_mode=True)
                
                if isinstance(response, dict):
                    terms_found = response.get("terms_found", [])
                    
                    for term_info in terms_found:
                        term = term_info.get("term", "")
                        found = term_info.get("found", False)
                        confidence = term_info.get("confidence", 0)
                        location = term_info.get("location", "")
                        
                        if found and confidence >= 0.6:
                            resultats["obligatoires"]["conformes"].append({
                                "terme": term,
                                "confidence": confidence,
                                "location": location
                            })
                        elif term in required_terms_23:
                            resultats["obligatoires"]["absents"].append({
                                "terme": term,
                                "raison": f"Confiance trop faible ({confidence:.2f}) ou non trouvé"
                            })
                            
            except Exception as e:
                print(f"⚠️ Erreur vérification glossaires LLM: {e}")
                # Fallback: recherche textuelle simple
                glossary_text_lower = glossary_text.lower()
                for term in required_terms_23:
                    term_lower = term.lower()
                    if term_lower and term_lower in glossary_text_lower:
                        resultats["obligatoires"]["conformes"].append({
                            "terme": term,
                            "confidence": 0.9,
                            "location": "Glossaire"
                        })
                    else:
                        resultats["obligatoires"]["absents"].append({
                            "terme": term,
                            "raison": "Non trouvé dans le glossaire"
                        })
        else:
            # Si pas de texte glossaire, tous les termes sont absents
            for term in required_terms_23:
                resultats["obligatoires"]["absents"].append({
                    "terme": term,
                    "raison": "Pas de section glossaire détectée"
                })
        
        # Vérification de l'utilisation des glossaires dans le document
        for glossaire in glossaires_attendus:
            if not glossaire:
                continue
            
            found = False
            glossaire_lower = glossaire.lower()
            
            # Vérifier d'abord si c'est un terme obligatoire déjà trouvé
            if any(glossaire.lower() == item.get("terme", "").lower() 
                   for item in resultats["obligatoires"]["conformes"]):
                # C'est un terme obligatoire conforme, ne pas le mettre dans utilisation
                continue
            
            # Rechercher dans tout le document
            for slide_data in state["slides_data"]:
                if glossaire_lower in slide_data["text"].lower():
                    found = True
                    self.metrics.glossaires_detectes += 1
                    
                    resultats["utilisation"]["conformes"].append({
                        "terme": glossaire,
                        "slide": slide_data["index"] + 1,
                        "contexte": slide_data["text"][:200] + "..."
                    })
                    break
            
            if not found:
                resultats["utilisation"]["absents"].append({
                    "terme": glossaire,
                    "raison": "Non trouvé dans le document"
                })
        
        # Calcul des scores
        total_required = len(required_terms_23)
        obligatoire_score = len(resultats["obligatoires"]["conformes"]) / max(total_required, 1)
        utilisation_score = len(resultats["utilisation"]["conformes"]) / max(len(glossaires_attendus), 1)
        
        # Créer un score global intelligent
        if glossaires_attendus:
            global_score = (obligatoire_score * 0.6 + utilisation_score * 0.4)
        else:
            global_score = obligatoire_score
        
        state["glossary_check"] = {
            **resultats,
            "obligatoire_score": obligatoire_score,
            "utilisation_score": utilisation_score,
            "global_score": global_score,
            "total_required": total_required,
            "found_count": len(resultats["obligatoires"]["conformes"]),
            "glossary_slide": glossary_slide_idx + 1 if glossary_slide_idx != -1 else "Non trouvé",
            "glossary_text_length": len(glossary_text),
            "note_explication": "Les termes dans 'obligatoires_conformes' sont les 23 glossaires obligatoires détectés. Les termes dans 'utilisation_conformes' sont les glossaires spécifiques au document détectés dans le contenu."
        }
        
        self.metrics.add_coverage_score(obligatoire_score)
        
        return state

    def verify_sources_intelligently_node(self, state: CombinedAgentState) -> CombinedAgentState:
        """Vérification intelligente des sources - OPTIMISÉ"""
        print("📚 Vérification intelligente des sources...")
        
        resultats = {
            "conformes": [], 
            "non_conformes": [], 
            "absents": [],
            "note": "Sources vérifiées pour leur présence et pertinence"
        }
        
        # Collecter toutes les sources des slides
        slides_with_sources = []
        for idx, slide_data in enumerate(state["slides_data"]):
            source_texte = slide_data["footer"]["source"]["texte"]
            
            if source_texte and not slide_data["footer"]["source"]["vide"]:
                slides_with_sources.append({
                    "index": idx,
                    "source": source_texte,
                    "slide_content": slide_data["text"][:300]  # Limiter la taille
                })
                self.metrics.sources_verifiees += 1
        
        print(f"  📊 {len(slides_with_sources)} slides avec des sources détectées")
        
        # STRATÉGIE OPTIMISÉE: Analyser les sources par groupes
        if slides_with_sources:
            for i in range(0, len(slides_with_sources), 5):
                group = slides_with_sources[i:i+5]
                
                system_prompt = """Évalue la qualité des sources dans un document financier.
                Réponds UNIQUEMENT en JSON:
                {
                    "sources_evaluated": [
                        {
                            "slide_index": numéro,
                            "quality_score": 0.0-1.0,
                            "valid_source": true/false,
                            "reason": "brief explanation",
                            "suggested_improvement": "suggestion d'amélioration si nécessaire"
                        }
                    ],
                    "summary": {
                        "total_evaluated": nombre,
                        "average_quality": moyenne,
                        "compliant_sources": nombre
                    }
                }"""
                
                user_prompt = f"""Évalue la qualité des sources suivantes:

{json.dumps(group, ensure_ascii=False, indent=2)}

Pour chaque source:
1. Score de qualité (0-1) - basé sur la clarté, complétude, et pertinence
2. Source valide (true/false) - est-ce une référence crédible?
3. Justification brève
4. Suggestion d'amélioration si le score < 0.6

Critères de qualité:
- Source clairement identifiable
- Référence complète (auteur, date, organisation si applicable)
- Pertinente par rapport au contenu
- Format cohérent"""

                try:
                    response, _ = self._call_llm(user_prompt, system_prompt, temperature=0.1, json_mode=True)
                    
                    if isinstance(response, dict):
                        sources_evaluated = response.get("sources_evaluated", [])
                        
                        for source_eval in sources_evaluated:
                            idx = source_eval.get("slide_index")
                            quality_score = source_eval.get("quality_score", 0)
                            valid = source_eval.get("valid_source", False)
                            reason = source_eval.get("reason", "")
                            
                            if valid and quality_score >= 0.6:
                                resultats["conformes"].append({
                                    "slide": idx + 1,
                                    "quality_score": round(quality_score, 2),
                                    "reason": reason,
                                    "source_text": group[i]["source"] if i < len(group) else ""
                                })
                            else:
                                resultats["non_conformes"].append({
                                    "slide": idx + 1,
                                    "quality_score": round(quality_score, 2),
                                    "raison": f"{reason} (score: {quality_score:.2f})",
                                    "suggestion": source_eval.get("suggested_improvement", "")
                                })
                except Exception as e:
                    print(f"⚠️ Erreur analyse sources groupe: {e}")
                    # Fallback: validation basique
                    for slide in group:
                        source_text = slide["source"]
                        if len(source_text) > 10 and not re.match(r'^\d{4}$', source_text.strip()):
                            resultats["conformes"].append({
                                "slide": slide["index"] + 1,
                                "quality_score": 0.7,
                                "reason": "Source basique détectée",
                                "source_text": source_text
                            })
                        else:
                            resultats["non_conformes"].append({
                                "slide": slide["index"] + 1,
                                "raison": "Source trop courte ou seulement une date",
                                "source_text": source_text
                            })
        
        # Marquer les slides sans sources
        for idx, slide_data in enumerate(state["slides_data"]):
            source_texte = slide_data["footer"]["source"]["texte"]
            if not source_texte or slide_data["footer"]["source"]["vide"]:
                resultats["absents"].append({
                    "slide": idx + 1,
                    "raison": "Aucune source détectée dans le footer"
                })
        
        source_score = len(resultats["conformes"]) / max(len(state["slides_data"]), 1)
        
        state["source_check"] = {
            **resultats,
            "score": source_score,
            "total_slides": len(state["slides_data"]),
            "slides_with_sources": len(slides_with_sources)
        }
        
        self.metrics.add_compliance_score(source_score)
        
        return state

    def verify_registration_intelligently_node(self, state: CombinedAgentState) -> CombinedAgentState:
        """Vérification intelligente de l'enregistrement du fonds avec gestion multi-ISIN"""
        print("📋 Vérification intelligente de l'enregistrement...")
        
        ref_fund_data = state.get("ref_fund_data", {})
        all_registered_countries = state.get("all_registered_countries", {})
        detected_countries = state.get("detected_countries", [])
        normalized_countries = state.get("normalized_countries", [])
        detected_isins = state.get("detected_fund", {}).get("isins", [])
        
        results = {
            "pays_corrects": [],
            "pays_manquants": [],
            "pays_exces": [],
            "isins_trouves": [],
            "isins_non_trouves": [],
            "details": {}
        }
        
        # Vérifier les ISINs trouvés
        if detected_isins:
            resources = self._load_resources()
            all_funds = resources.get('fonds', [])
            
            for isin in detected_isins:
                found = False
                for fund in all_funds:
                    if fund.get('isin') == isin:
                        found = True
                        results["isins_trouves"].append({
                            "isin": isin,
                            "fonds": fund.get('nom_fonds', 'Inconnu')
                        })
                        break
                
                if not found:
                    results["isins_non_trouves"].append({
                        "isin": isin,
                        "raison": "Non trouvé dans la base de données"
                    })
        
        if not ref_fund_data and not all_registered_countries:
            state["registration_check"] = {
                "status": "FONDS_NON_IDENTIFIE",
                "reason": "Fonds non trouvé dans la base de données",
                "isins": detected_isins,
                "results": results,
                "compliant": False,
                "score": 0
            }
            return state
        
        # Liste des pays enregistrés (normalisés)
        registered_countries_list = list(all_registered_countries.keys())
        
        print(f"  🌍 Comparaison des pays:")
        print(f"     Détectés: {detected_countries}")
        print(f"     Normalisés: {normalized_countries}")
        print(f"     Enregistrés: {registered_countries_list}")
        
        # Comparaison des pays détectés avec les pays enregistrés
        for detected_country in detected_countries:
            normalized_country = self.country_detector.normalize_country_name(detected_country)
            if not normalized_country:
                continue
                
            found = False
            for registered_country in registered_countries_list:
                if normalized_country.lower() == registered_country.lower():
                    found = True
                    reg_info = all_registered_countries.get(registered_country, {})
                    status_code = reg_info.get('info', {}).get('code', '') if isinstance(reg_info, dict) else ''
                    
                    if status_code in ['R', 'RP']:
                        results["pays_corrects"].append({
                            "pays_document": detected_country,
                            "pays_normalise": normalized_country,
                            "pays_base": registered_country,
                            "status": status_code,
                            "compliant": True
                        })
                    else:
                        results["pays_manquants"].append({
                            "pays_document": detected_country,
                            "pays_normalise": normalized_country,
                            "pays_base": registered_country,
                            "status": status_code,
                            "raison": "Pays mentionné mais non enregistré (R/RP)"
                        })
                    break
            
            if not found:
                results["pays_exces"].append({
                    "pays_document": detected_country,
                    "pays_normalise": normalized_country,
                    "raison": "Pays mentionné mais non trouvé dans l'enregistrement du fonds"
                })
        
        # Vérifier si des pays enregistrés ne sont pas mentionnés
        for registered_country, reg_info in all_registered_countries.items():
            status_code = reg_info.get('info', {}).get('code', '') if isinstance(reg_info, dict) else ''
            
            if status_code not in ['R', 'RP']:
                continue
                
            found = False
            for normalized_country in normalized_countries:
                if normalized_country.lower() == registered_country.lower():
                    found = True
                    break
            
            if not found:
                results["pays_manquants"].append({
                    "pays_base": registered_country,
                    "pays_original": reg_info.get('original', registered_country),
                    "status": status_code,
                    "raison": "Pays enregistré mais non mentionné dans le document"
                })
        
        # Calcul du score
        total_checks = len(normalized_countries) + len([c for c in all_registered_countries.keys() 
                                                       if all_registered_countries.get(c, {}).get('info', {}).get('code', '') in ['R', 'RP']])
        correct_checks = len(results["pays_corrects"])
        
        score = correct_checks / max(total_checks, 1) if total_checks > 0 else 1.0
        
        # Calculer un score pour les ISINs
        isin_score = len(results["isins_trouves"]) / max(len(detected_isins), 1) if detected_isins else 1.0
        
        # Score global combiné
        global_score = (score * 0.7 + isin_score * 0.3) if detected_isins else score
        
        state["registration_check"] = {
            "fonds_identifie": ref_fund_data.get("nom_fonds", "Non trouvé"),
            "isins": detected_isins,
            "pays_detectes": detected_countries,
            "pays_normalises": normalized_countries,
            "pays_enregistrement": all_registered_countries,
            "results": results,
            "score": global_score,
            "pays_score": score,
            "isin_score": isin_score,
            "compliant": len(results["pays_manquants"]) == 0 and len(results["pays_exces"]) == 0,
            "summary": {
                "pays_corrects": len(results["pays_corrects"]),
                "pays_manquants": len(results["pays_manquants"]),
                "pays_exces": len(results["pays_exces"]),
                "isins_trouves": len(results["isins_trouves"]),
                "isins_non_trouves": len(results["isins_non_trouves"])
            },
            "note": "Vérification basée sur la correspondance sémantique des pays et la présence des ISINs"
        }
        
        self.metrics.add_compliance_score(global_score)
        
        return state

    def extract_fund_characteristics_node(self, state: CombinedAgentState) -> CombinedAgentState:
        """Extraction intelligente des caractéristiques du fonds"""
        print("📊 Extraction intelligente des caractéristiques...")
        
        characteristics_slide_idx = state.get("characteristics_slide_idx", -1)
        
        # Si le slide des caractéristiques a été identifié
        if characteristics_slide_idx != -1 and characteristics_slide_idx < len(state["slides_data"]):
            slide_data = state["slides_data"][characteristics_slide_idx]
            slide_text = slide_data["text"]
            
            system_prompt = """Tu es un expert en analyse de caractéristiques de fonds financiers.
            Extrait les caractéristiques du fonds à partir du texte.
            Réponds UNIQUEMENT en JSON:
            {
                "caracteristiques_trouvees": {
                    "isin": "valeur ou null",
                    "devise": "valeur ou null",
                    "valorisation": "valeur ou null",
                    "frais": "valeur ou null",
                    "performance": "valeur ou null",
                    "risque": "valeur ou null",
                    "classes_parts": "valeur ou null",
                    "horizon_investissement": "valeur ou null",
                    "strategie": "valeur ou null"
                },
                "completeness_score": 0.0-1.0,
                "found_count": nombre,
                "missing_elements": ["éléments manquants"],
                "note": "note sur l'extraction"
            }"""
            
            user_prompt = f"""Extrait les caractéristiques du fonds à partir du texte suivant:

{slide_text}

Identifie les informations suivantes si présentes:
- ISIN (doit correspondre à un format valide)
- Devise de base (EUR, USD, CHF, etc.)
- Fréquence de valorisation (quotidienne, mensuelle, etc.)
- Frais (de gestion, d'entrée, de sortie, maximum)
- Performance historique ou objectif de performance
- Niveau de risque (1-7 ou description)
- Classes de parts disponibles
- Horizon d'investissement recommandé
- Stratégie d'investissement"""
            
            try:
                response, _ = self._call_llm(user_prompt, system_prompt, temperature=0.1, json_mode=True)
                
                if isinstance(response, dict):
                    caracteristiques = response.get("caracteristiques_trouvees", {})
                    completeness_score = response.get("completeness_score", 0)
                    found_count = response.get("found_count", 0)
                    missing = response.get("missing_elements", [])
                    
                    state["characteristics_check"] = {
                        "slide": characteristics_slide_idx + 1,
                        "caracteristiques": caracteristiques,
                        "found_count": found_count,
                        "total_expected": 9,  # Nombre d'éléments recherchés
                        "missing_elements": missing,
                        "completeness_score": round(completeness_score, 3),
                        "autres_caracteristiques": [],
                        "status": "CONFORME" if completeness_score >= 0.5 else "PARTIELLEMENT CONFORME",
                        "note": response.get("note", "")
                    }
                else:
                    state["characteristics_check"] = {
                        "slide": characteristics_slide_idx + 1,
                        "caracteristiques": {},
                        "found_count": 0,
                        "total_expected": 9,
                        "completeness_score": 0,
                        "status": "NON CONFORME",
                        "raison": "Erreur d'extraction LLM"
                    }
                    
            except Exception as e:
                print(f"⚠️ Erreur extraction caractéristiques: {e}")
                state["characteristics_check"] = {
                    "slide": characteristics_slide_idx + 1,
                    "caracteristiques": {},
                    "found_count": 0,
                    "total_expected": 9,
                    "completeness_score": 0,
                    "status": "NON CONFORME",
                    "raison": f"Erreur: {str(e)}"
                }
        else:
            # Recherche dans tout le document en un seul appel
            system_prompt = """Tu es un expert en analyse de caractéristiques de fonds financiers.
            Extrait les caractéristiques du fonds à partir du texte complet.
            Réponds UNIQUEMENT en JSON avec la même structure."""
            
            user_prompt = f"""Extrait les caractéristiques du fonds à partir du document complet:

{state['full_text'][:4000]}... [document tronqué]

Identifie les informations importantes du fonds."""
            
            try:
                response, _ = self._call_llm(user_prompt, system_prompt, temperature=0.1, json_mode=True)
                
                if isinstance(response, dict):
                    caracteristiques = response.get("caracteristiques_trouvees", {})
                    completeness_score = response.get("completeness_score", 0)
                    found_count = response.get("found_count", 0)
                    missing = response.get("missing_elements", [])
                    
                    state["characteristics_check"] = {
                        "slide": "Document entier",
                        "caracteristiques": caracteristiques,
                        "found_count": found_count,
                        "completeness_score": round(completeness_score, 3),
                        "missing_elements": missing,
                        "status": "PARTIELLEMENT CONFORME" if completeness_score >= 0.3 else "NON CONFORME",
                        "note": response.get("note", "")
                    }
                else:
                    state["characteristics_check"] = {
                        "slide": None,
                        "caracteristiques": {},
                        "completeness_score": 0,
                        "status": "NON CONFORME",
                        "raison": "Caractéristiques non trouvées"
                    }
                    
            except Exception as e:
                print(f"⚠️ Erreur extraction caractéristiques globale: {e}")
                state["characteristics_check"] = {
                    "slide": None,
                    "caracteristiques": {},
                    "completeness_score": 0,
                    "status": "NON CONFORME",
                    "raison": "Erreur d'analyse"
                }
        
        self.metrics.add_coverage_score(state["characteristics_check"].get("completeness_score", 0))
        
        return state

    def generate_comprehensive_report_node(self, state: CombinedAgentState) -> CombinedAgentState:
        """Génération d'un rapport complet et détaillé avec information LLM"""
        print("📝 Génération du rapport complet...")
        
        self.metrics.stop()
        
        # Ajouter l'information sur les appels LLM
        llm_warning = None
        if self.llm_call_count > self.max_llm_calls:
            llm_warning = f"⚠️ INFORMATION: {self.llm_call_count} appels LLM effectués (seuil de référence: {self.max_llm_calls}). L'analyse a été complète sans limitation de performance."
            print(llm_warning)
        
        # Calcul des scores globaux
        disclaimer_score = state.get("disclaimer_check", {}).get("global_score", 0)
        glossary_score = state.get("glossary_check", {}).get("global_score", 0)
        source_score = state.get("source_check", {}).get("score", 0)
        registration_score = state.get("registration_check", {}).get("score", 0)
        characteristics_score = state.get("characteristics_check", {}).get("completeness_score", 0)
        
        scores = [s for s in [disclaimer_score, glossary_score, source_score, 
                             registration_score, characteristics_score] if s > 0]
        global_compliance = np.mean(scores) if scores else 0
        
        # Nom du fichier de output
        pptx_filename = os.path.basename(state["pptx_path"])
        base_name = os.path.splitext(pptx_filename)[0]
        output_dir = "workflow/output/rapport"
        
        # Créer le dossier output s'il n'existe pas
        os.makedirs(output_dir, exist_ok=True)
        
        output_filename = f"{base_name}_dis_glos.json"
        output_path = os.path.join(output_dir, output_filename)
        
        # Récupérer les résultats détaillés
        registration_check = state.get("registration_check", {})
        results = registration_check.get("results", {}) if registration_check else {}
        
        # S'assurer que tous les champs existent
        disclaimer_check = state.get("disclaimer_check", {})
        glossary_check = state.get("glossary_check", {})
        source_check = state.get("source_check", {})
        characteristics_check = state.get("characteristics_check", {})
        
        # Données de référence
        ref_fund_data = state.get("ref_fund_data", {})
        all_registered_countries = state.get("all_registered_countries", {})
        
        report = {
            "timestamp": datetime.now().isoformat(),
            "document": pptx_filename,
            "glossaires_source": os.path.basename(state["glossaires_json_path"]),
            "information_llm": {
                "appels_total": self.llm_call_count,
                "seuil_reference": self.max_llm_calls,
                "warning": llm_warning,
                "note": f"L'analyse a utilisé {self.llm_call_count} appels LLM pour garantir une analyse complète et précise."
            },
            
            "statistiques_globales": {
                "total_slides": len(state["slides_data"]),
                "langue_detectee": state.get("detected_lang", "fr"),
                "cible_detectee": state.get("detected_target", "Professionnel"),
                "fonds_identifie": state.get("detected_fund", {}).get("name", "Inconnu"),
                "isins_detectes": state.get("detected_fund", {}).get("isins", []),
                "classe_parts": state.get("detected_fund", {}).get("classe_part", "Non spécifiée"),
                "pays_commercialisation_detectes": state.get("detected_countries", []),
                "pays_commercialisation_normalises": state.get("normalized_countries", []),
                "score_global_conformite": round(global_compliance, 3),
                "niveau_conformite": self._get_compliance_level(global_compliance)
            },
            
            "resultats_detailles": {
                "enregistrement": {
                    "fonds_reference": ref_fund_data.get("nom_fonds", "Non trouvé"),
                    "isins_reference": ref_fund_data.get("isin", "Non trouvé"),
                    "pays_enregistrement": all_registered_countries,
                    "pays_corrects": results.get("pays_corrects", []),
                    "pays_manquants": results.get("pays_manquants", []),
                    "pays_exces": results.get("pays_exces", []),
                    "isins_trouves": results.get("isins_trouves", []),
                    "isins_non_trouves": results.get("isins_non_trouves", []),
                    "score": registration_check.get("score", 0),
                    "pays_score": registration_check.get("pays_score", 0),
                    "isin_score": registration_check.get("isin_score", 0),
                    "compliant": registration_check.get("compliant", False),
                    "summary": registration_check.get("summary", {})
                },
                
                "disclaimers": disclaimer_check,
                
                "glossaires": glossary_check,
                
                "sources": source_check,
                
                "caracteristiques": characteristics_check
            },
            
            "scores_individuels": {
                "enregistrement": round(registration_score, 3),
                "disclaimers": round(disclaimer_score, 3),
                "glossaires": round(glossary_score, 3),
                "sources": round(source_score, 3),
                "caracteristiques": round(characteristics_score, 3)
            },
            
            "metriques_performance": self.metrics.to_dict(),
            "erreurs": state.get("erreurs", []),
            "recommandations": self._generate_intelligent_recommendations(state),
            "resume_executif": self._generate_executive_summary(state, global_compliance),
            
            "metadonnees": {
                "agent_version": "7.0_optimisee",
                "llm_calls_count": self.llm_call_count,
                "llm_calls_threshold": self.max_llm_calls,
                "llm_client_metrics": self.llm_client.get_metrics(),
                "date_analyse": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                "optimisations": {
                    "detection_isin_ameliorée": True,
                    "detection_pays_robuste": True,
                    "validation_isin_stricte": True,
                    "grouped_llm_calls": True,
                    "text_truncation": True,
                    "fallback_strategies": True,
                    "performance_monitoring": True
                }
            }
        }
        
        # Sauvegarde du rapport
        try:
            with open(output_path, 'w', encoding='utf-8') as f:
                json.dump(report, f, indent=2, ensure_ascii=False)
            
            state["report"] = report
            state["llm_warning"] = llm_warning
            
            # Affichage du résumé
            self._print_comprehensive_summary(report)
            
            print(f"\n✅ Rapport sauvegardé: {output_path}")
            
        except Exception as e:
            print(f"❌ Erreur sauvegarde rapport: {e}")
            state["erreurs"].append(f"Erreur sauvegarde rapport: {e}")
        
        return state

    def _get_compliance_level(self, score: float) -> str:
        """Détermine le niveau de conformité"""
        if score >= 0.9:
            return "EXCELLENT"
        elif score >= 0.7:
            return "BON"
        elif score >= 0.5:
            return "MOYEN"
        elif score >= 0.3:
            return "FAIBLE"
        else:
            return "NON CONFORME"

    def _generate_intelligent_recommendations(self, state: CombinedAgentState) -> List[str]:
        """Génère des recommandations intelligentes basées sur les résultats"""
        recommandations = []
        
        # Glossaires
        glossary_check = state.get("glossary_check", {})
        obligatoires_absents = glossary_check.get("obligatoires", {}).get("absents", [])
        if obligatoires_absents:
            count = len(obligatoires_absents)
            termes = [item.get("terme", "") for item in obligatoires_absents[:3]]
            recommandations.append(f"Ajouter {count} glossaires obligatoires manquants parmi les 23 requis. Exemples: {', '.join(termes)}")
        
        utilisation_absents = glossary_check.get("utilisation", {}).get("absents", [])
        if utilisation_absents:
            count = len(utilisation_absents)
            recommandations.append(f"Vérifier l'utilisation des {count} glossaires spécifiques manquants dans le document")
        
        # Disclaimers
        disclaimer_check = state.get("disclaimer_check", {})
        obligatoires_absents = disclaimer_check.get("obligatoires", {}).get("absents", [])
        if obligatoires_absents:
            count = len(obligatoires_absents)
            recommandations.append(f"Ajouter {count} disclaimers obligatoires manquants pour la langue {state.get('detected_lang')}")
        
        # Alignement des disclaimers
        alignement_non_conformes = disclaimer_check.get("alignement", {}).get("non_conformes", [])
        if alignement_non_conformes:
            slides = [str(item.get('slide', '')) for item in alignement_non_conformes[:5]]
            recommandations.append(f"Améliorer l'alignement des disclaimers sur les slides: {', '.join(slides)}")
        
        # Sources
        source_check = state.get("source_check", {})
        if "absents" in source_check:
            absent_count = len(source_check["absents"])
            if absent_count > 0:
                recommandations.append(f"Ajouter des sources sur {absent_count} slides sans référence")
        
        # Enregistrement
        registration_check = state.get("registration_check", {})
        results = registration_check.get("results", {})
        pays_manquants = results.get("pays_manquants", [])
        pays_exces = results.get("pays_exces", [])
        
        if pays_manquants:
            pays_list = [p.get('pays_base', p.get('pays_document', p.get('pays_normalise', ''))) for p in pays_manquants[:3]]
            recommandations.append(f"Vérifier l'enregistrement pour {len(pays_manquants)} pays: {', '.join(pays_list)}")
        
        if pays_exces:
            pays_list = [p.get('pays_document', p.get('pays_normalise', '')) for p in pays_exces[:3]]
            recommandations.append(f"Vérifier les mentions des {len(pays_exces)} pays non enregistrés: {', '.join(pays_list)}")
        
        # Caractéristiques
        characteristics_check = state.get("characteristics_check", {})
        if characteristics_check.get("completeness_score", 0) < 0.5:
            missing = characteristics_check.get("missing_elements", [])
            if missing:
                recommandations.append(f"Compléter les caractéristiques du fonds: {', '.join(missing[:3])}")
            else:
                recommandations.append("Améliorer la présentation des caractéristiques du fonds (ISIN, devise, frais, risque, etc.)")
        
        # ISINs non trouvés
        isins_non_trouves = results.get("isins_non_trouves", [])
        if isins_non_trouves:
            isin_list = [i.get('isin', '') for i in isins_non_trouves[:3]]
            recommandations.append(f"Vérifier les {len(isins_non_trouves)} ISINs non trouvés dans la base: {', '.join(isin_list)}")
        
        # Ajouter une recommandation sur les appels LLM si nécessaire
        if self.llm_call_count > self.max_llm_calls:
            recommandations.append(f"Optimisation: L'analyse a utilisé {self.llm_call_count} appels LLM. Pour les documents standards, des optimisations peuvent réduire ce nombre sans affecter la qualité.")
        
        return recommandations

    def _generate_executive_summary(self, state: CombinedAgentState, global_score: float) -> Dict:
        """Génère un résumé exécutif intelligent"""
        total_slides = len(state["slides_data"])
        
        summary = {
            "conformite_globale": f"{global_score:.1%}",
            "niveau": self._get_compliance_level(global_score),
            "points_forts": [],
            "points_faibles": [],
            "actions_prioritaires": [],
            "analyse_llm": {
                "appels": self.llm_call_count,
                "note": f"Analyse complète avec {self.llm_call_count} appels LLM" +
                       (f" (dépassement du seuil de {self.max_llm_calls})" if self.llm_call_count > self.max_llm_calls else "")
            }
        }
        
        # Points forts
        if state.get("glossary_check", {}).get("obligatoire_score", 0) >= 0.8:
            summary["points_forts"].append("Glossaires obligatoires bien couverts")
        
        if state.get("disclaimer_check", {}).get("obligatoire_score", 0) >= 0.8:
            summary["points_forts"].append("Disclaimers obligatoires présents")
        
        registration_check = state.get("registration_check", {})
        if registration_check.get("score", 0) >= 0.8:
            summary["points_forts"].append("Enregistrement du fonds globalement conforme")
        
        # Points faibles
        if state.get("source_check", {}).get("score", 0) < 0.5:
            summary["points_faibles"].append("Sources insuffisantes ou absentes")
        
        if state.get("characteristics_check", {}).get("completeness_score", 0) < 0.5:
            summary["points_faibles"].append("Caractéristiques du fonds incomplètes")
        
        registration_results = registration_check.get("results", {})
        if len(registration_results.get("pays_manquants", [])) > 0:
            summary["points_faibles"].append(f"{len(registration_results['pays_manquants'])} problèmes d'enregistrement de pays")
        
        if len(registration_results.get("pays_exces", [])) > 0:
            summary["points_faibles"].append(f"{len(registration_results['pays_exces'])} pays mentionnés non enregistrés")
        
        # Actions prioritaires
        recommandations = self._generate_intelligent_recommendations(state)
        summary["actions_prioritaires"] = recommandations[:3]  # Top 3 recommandations
        
        return summary

    def _print_comprehensive_summary(self, report: Dict):
        """Affiche un résumé complet du rapport"""
        print("\n" + "="*80)
        print("RAPPORT DE CONFORMITÉ OPTIMISÉ - RÉSUMÉ EXÉCUTIF")
        print("="*80)
        
        stats = report["statistiques_globales"]
        scores = report["scores_individuels"]
        summary = report["resume_executif"]
        llm_info = report.get("information_llm", {})
        
        print(f"\n📊 CONFORMITÉ GLOBALE: {summary['conformite_globale']} ({summary['niveau']})")
        print(f"   Document: {report['document']}")
        print(f"   Fonds: {stats['fonds_identifie']}")
        print(f"   ISINs valides détectés: {', '.join(stats.get('isins_detectes', [])) if stats.get('isins_detectes') else 'Aucun'}")
        print(f"   Slides: {stats['total_slides']} | Langue: {stats['langue_detectee']} | Cible: {stats['cible_detectee']}")
        
        # Information LLM
        if llm_info.get("warning"):
            print(f"\n⚠️  {llm_info['warning']}")
        
        print(f"\n📈 SCORES PAR CATÉGORIE:")
        print(f"   • Enregistrement: {scores['enregistrement']:.1%}")
        print(f"   • Disclaimers: {scores['disclaimers']:.1%}")
        print(f"   • Glossaires: {scores['glossaires']:.1%}")
        print(f"   • Sources: {scores['sources']:.1%}")
        print(f"   • Caractéristiques: {scores['caracteristiques']:.1%}")
        
        # Détails de l'enregistrement
        registration = report["resultats_detailles"]["enregistrement"]
        if registration.get("pays_corrects") or registration.get("pays_manquants") or registration.get("pays_exces"):
            print(f"\n🌍 ENREGISTREMENT PAYS:")
            
            if registration["pays_corrects"]:
                print(f"   ✅ Corrects: {len(registration['pays_corrects'])} pays")
                for i, pays in enumerate(registration["pays_corrects"][:3]):
                    nom = pays.get('pays_document', pays.get('pays_normalise', ''))
                    status = pays.get('status', '')
                    print(f"      {i+1}. {nom} ({status})")
                if len(registration["pays_corrects"]) > 3:
                    print(f"      ... et {len(registration['pays_corrects']) - 3} autres")
            
            if registration["pays_manquants"]:
                print(f"   ❌ Manquants: {len(registration['pays_manquants'])} pays")
                for i, pays in enumerate(registration["pays_manquants"][:3]):
                    nom = pays.get('pays_base', pays.get('pays_document', ''))
                    raison = pays.get('raison', '')
                    print(f"      {i+1}. {nom} ({raison})")
            
            if registration["pays_exces"]:
                print(f"   ⚠️  Excès: {len(registration['pays_exces'])} pays")
                for i, pays in enumerate(registration["pays_exces"][:3]):
                    nom = pays.get('pays_document', pays.get('pays_normalise', ''))
                    print(f"      {i+1}. {nom}")
        
        # ISINs
        if registration.get("isins_trouves") or registration.get("isins_non_trouves"):
            print(f"\n🔢 ISINS:")
            
            if registration["isins_trouves"]:
                print(f"   ✅ Trouvés: {len(registration['isins_trouves'])} ISIN(s)")
                for i, isin in enumerate(registration["isins_trouves"][:3]):
                    print(f"      {i+1}. {isin.get('isin')} - {isin.get('fonds')}")
            
            if registration["isins_non_trouves"]:
                print(f"   ❌ Non trouvés: {len(registration['isins_non_trouves'])} ISIN(s)")
                for i, isin in enumerate(registration["isins_non_trouves"][:3]):
                    print(f"      {i+1}. {isin.get('isin')}")
        
        print(f"\n✅ POINTS FORTS:")
        for point in summary.get("points_forts", []):
            print(f"   ✓ {point}")
        
        print(f"\n⚠️  POINTS FAIBLES:")
        for point in summary.get("points_faibles", []):
            print(f"   ✗ {point}")
        
        print(f"\n🔧 ACTIONS PRIORITAIRES:")
        for rec in summary.get("actions_prioritaires", []):
            print(f"   • {rec}")
        
        print(f"\n⏱️  PERFORMANCE:")
        perf = report["metriques_performance"]
        print(f"   Durée totale: {perf['duree_execution_secondes']}s")
        print(f"   Appels LLM: {perf['appels_llm']} (total: {llm_info.get('appels_total', 0)})")
        print(f"   Slides/s: {perf['efficacite']['slides_par_seconde']:.2f}")
        print(f"   Taux succès: {perf['efficacite']['taux_succes']:.1%}")
        
        print("="*80)

    def build_graph(self):
        """Construit le graphe LangGraph optimisé"""
        workflow = StateGraph(CombinedAgentState)
        
        workflow.add_node("parse_document", self.parse_document_node)
        workflow.add_node("load_glossaires", self.load_glossaires_node)
        workflow.add_node("identify_fund_and_context", self.identify_fund_and_context_node)
        workflow.add_node("verify_disclaimers_intelligently", self.verify_disclaimers_intelligently_node)
        workflow.add_node("verify_glossaries_intelligently", self.verify_glossaries_intelligently_node)
        workflow.add_node("verify_sources_intelligently", self.verify_sources_intelligently_node)
        workflow.add_node("verify_registration_intelligently", self.verify_registration_intelligently_node)
        workflow.add_node("extract_fund_characteristics", self.extract_fund_characteristics_node)
        workflow.add_node("generate_comprehensive_report", self.generate_comprehensive_report_node)
        
        workflow.set_entry_point("parse_document")
        workflow.add_edge("parse_document", "load_glossaires")
        workflow.add_edge("load_glossaires", "identify_fund_and_context")
        workflow.add_edge("identify_fund_and_context", "verify_disclaimers_intelligently")
        workflow.add_edge("verify_disclaimers_intelligently", "verify_glossaries_intelligently")
        workflow.add_edge("verify_glossaries_intelligently", "verify_sources_intelligently")
        workflow.add_edge("verify_sources_intelligently", "verify_registration_intelligently")
        workflow.add_edge("verify_registration_intelligently", "extract_fund_characteristics")
        workflow.add_edge("extract_fund_characteristics", "generate_comprehensive_report")
        workflow.add_edge("generate_comprehensive_report", END)
        
        return workflow.compile()

    def run(self, pptx_path: str, glossaires_json_path: str):
        """Exécute l'agent robuste sur un document PPTX"""
        if not os.path.exists(pptx_path):
            print(f"❌ Erreur: Le fichier PPTX '{pptx_path}' n'existe pas.")
            return None
        
        if not os.path.exists(glossaires_json_path):
            print(f"❌ Erreur: Le fichier JSON '{glossaires_json_path}' n'existe pas.")
            return None
        
        print("🚀 DÉMARRAGE DE L'ANALYSE OPTIMISÉE DE CONFORMITÉ")
        print("="*50)
        print(f"📄 Document: {os.path.basename(pptx_path)}")
        print(f"📖 Glossaires: {os.path.basename(glossaires_json_path)}")
        print(f"⚡ Version: 7.0_optimisee avec détection robuste")
        print("="*50)
        
        try:
            graph = self.build_graph()
            initial_state = {
                "pptx_path": pptx_path,
                "glossaires_json_path": glossaires_json_path,
                "erreurs": [],
                "detected_countries": [],
                "normalized_countries": [],
                "current_slide_idx": 0,
                "annexe_slide_idx": -1,
                "characteristics_slide_idx": -1,
                "scores_alignement": [],
                "qualite_reponses": [],
                "llm_call_count": 0,
                "llm_warning": None
            }
            
            result = graph.invoke(initial_state)
            
            # Récupérer les métriques du client LLM
            llm_metrics = self.llm_client.get_metrics()
            print(f"\n📊 MÉTRIQUES LLM:")
            print(f"   • Appels totaux: {llm_metrics['total_calls']}")
            print(f"   • Temps total: {llm_metrics['total_time_seconds']:.1f}s")
            print(f"   • Taux succès: {llm_metrics['success_rate']:.1f}%")
            print(f"   • Appels optimisés: {self.llm_call_count} (seuil: {self.max_llm_calls})")
            
            if self.llm_call_count > self.max_llm_calls:
                print(f"   ⚠️  Seuil dépassé de {self.llm_call_count - self.max_llm_calls} appels")
                print(f"   ℹ️  L'analyse reste complète sans limitation de performance")
            
            return result
            
        except Exception as e:
            print(f"\n❌ ERREUR CRITIQUE lors de l'exécution: {e}")
            import traceback
            traceback.print_exc()
            return None

# ==================== FONCTION PRINCIPALE ====================
def main():
    """Fonction principale pour l'exécution en ligne de commande"""
    # Créer le dossier output s'il n'existe pas
    os.makedirs("output/report", exist_ok=True)
    
    # Lancer l'agent optimisé
    agent = OptimizedDisclaimerGlossaryAgent(max_llm_calls=50)
    
    try:
        # Exemple d'utilisation - ajustez les chemins selon vos besoins
        pptx_path = "Batch 2/VERSION FINALE 6PG-GB-ODDO BHF Artificial Intelligence.pptx"
        glossaires_json_path = "workflow/caches/glossaires.json"
        
        result = agent.run(pptx_path, glossaires_json_path)
        
        if result:
            print("\n" + "="*60)
            print("✅ ANALYSE TERMINÉE AVEC SUCCÈS!")
            print("="*60)
            
            # Afficher le récapitulatif des appels LLM
            llm_count = agent.llm_call_count
            threshold = agent.max_llm_calls
            
            if llm_count > threshold:
                print(f"📊 APPELS LLM: {llm_count} (dépassement de {llm_count - threshold} appels)")
                print("   ℹ️  L'analyse a été complète sans limitation de performance")
                print("   ℹ️  Cette information est incluse dans le rapport")
            else:
                print(f"📊 APPELS LLM: {llm_count} (dans les limites du seuil de {threshold})")
            
            pptx_filename = os.path.basename(pptx_path)
            base_name = os.path.splitext(pptx_filename)[0]
            print(f"📄 Rapport disponible dans: output/{base_name}_dis_glos.json")
        else:
            print("\n❌ L'analyse a échoué.")
            
    except Exception as e:
        print(f"\n❌ ERREUR lors de l'exécution: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    main()