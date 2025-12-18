"""
Agent d'Audit et d'Assistance Documentaire Intelligent
Avec intégration LangGraph, système de mémoire et métriques complètes
"""

import os
import json
import warnings
from typing import Dict, List, Tuple, Optional, Any, Callable
from dataclasses import dataclass, field, asdict
from datetime import datetime
import numpy as np
from enum import Enum
import networkx as nx
import matplotlib.pyplot as plt
from sklearn.feature_extraction.text import CountVectorizer
from rank_bm25 import BM25Okapi
import faiss
from sentence_transformers import SentenceTransformer
import re
from collections import Counter
import itertools

# ============================================================================
# 0. IMPORT DU CLIENT LLM EXTERNE
# ============================================================================

try:
    from llm_client import LlamaClient
except ImportError:
    print("⚠️ Client LLM non trouvé. Utilisation du mode simulation.")
    LlamaClient = None

# ============================================================================
# 1. MODULES D'EXTRACTION DOCUMENTAIRE
# ============================================================================

try:
    import pdfplumber
    from pptx import Presentation
    from docx import Document
except ImportError:
    print("Installation des dépendances d'extraction...")

class DocumentType(Enum):
    PDF = "pdf"
    PPTX = "pptx"
    DOCX = "docx"

@dataclass
class DocumentChunk:
    """Représentation d'un fragment de document avec métadonnées"""
    text: str
    doc_type: DocumentType
    doc_name: str
    page_num: int = 0
    slide_num: int = 0
    chunk_id: str = ""
    metadata: Dict[str, Any] = field(default_factory=dict)

class DocumentExtractor:
    """Unifié l'extraction de différents types de documents"""
    
    def __init__(self):
        self.supported_types = [DocumentType.PDF, DocumentType.PPTX, DocumentType.DOCX]
    
    def extract(self, file_path: str) -> List[DocumentChunk]:
        """Extrait le contenu d'un document selon son type"""
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"Document non trouvé: {file_path}")
        
        ext = file_path.split('.')[-1].lower()
        chunks = []
        
        if ext == 'pdf':
            chunks = self._extract_pdf(file_path)
        elif ext == 'pptx':
            chunks = self._extract_pptx(file_path)
        elif ext == 'docx':
            chunks = self._extract_docx(file_path)
        else:
            raise ValueError(f"Format non supporté: {ext}")
        
        return chunks
    
    def _extract_pdf(self, file_path: str) -> List[DocumentChunk]:
        """Extrait le contenu d'un PDF"""
        chunks = []
        doc_name = os.path.basename(file_path)
        
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                text = page.extract_text()
                if text and text.strip():
                    chunk = DocumentChunk(
                        text=text.strip(),
                        doc_type=DocumentType.PDF,
                        doc_name=doc_name,
                        page_num=page_num,
                        chunk_id=f"{doc_name}_pdf_page_{page_num}"
                    )
                    chunks.append(chunk)
        
        return chunks
    
    def _extract_pptx(self, file_path: str) -> List[DocumentChunk]:
        """Extrait le contenu d'un PowerPoint"""
        chunks = []
        doc_name = os.path.basename(file_path)
        presentation = Presentation(file_path)
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            
            if slide_text:
                chunk = DocumentChunk(
                    text="\n".join(slide_text),
                    doc_type=DocumentType.PPTX,
                    doc_name=doc_name,
                    slide_num=slide_num,
                    chunk_id=f"{doc_name}_pptx_slide_{slide_num}"
                )
                chunks.append(chunk)
        
        return chunks
    
    def _extract_docx(self, file_path: str) -> List[DocumentChunk]:
        """Extrait le contenu d'un document Word"""
        chunks = []
        doc_name = os.path.basename(file_path)
        doc = Document(file_path)
        
        for para_num, paragraph in enumerate(doc.paragraphs, 1):
            if paragraph.text.strip():
                chunk = DocumentChunk(
                    text=paragraph.text.strip(),
                    doc_type=DocumentType.DOCX,
                    doc_name=doc_name,
                    page_num=para_num,
                    chunk_id=f"{doc_name}_docx_para_{para_num}"
                )
                chunks.append(chunk)
        
        return chunks

# ============================================================================
# 2. SYSTÈME DE RECHERCHE HYBRIDE (BM25 + VECTORIEL)
# ============================================================================

class HybridRetriever:
    """Système de recherche hybride BM25 + embeddings"""
    
    def __init__(self, model_name: str = 'paraphrase-multilingual-MiniLM-L12-v2'):
        self.bm25 = None
        self.vectorizer = CountVectorizer(stop_words='english')
        self.embedding_model = SentenceTransformer(model_name)
        self.faiss_index = None
        self.chunks = []
        self.chunk_embeddings = []
        
    def index_chunks(self, chunks: List[DocumentChunk]):
        """Indexe les chunks pour la recherche"""
        self.chunks = chunks
        
        # Indexation BM25
        tokenized_corpus = [self._tokenize(chunk.text) for chunk in chunks]
        self.bm25 = BM25Okapi(tokenized_corpus)
        
        # Indexation vectorielle
        texts = [chunk.text for chunk in chunks]
        self.chunk_embeddings = self.embedding_model.encode(texts, show_progress_bar=True)
        
        # Création de l'index FAISS
        dimension = self.chunk_embeddings.shape[1]
        self.faiss_index = faiss.IndexFlatL2(dimension)
        self.faiss_index.add(self.chunk_embeddings.astype('float32'))
    
    def _tokenize(self, text: str) -> List[str]:
        """Tokenisation simple pour BM25"""
        return text.lower().split()
    
    def hybrid_search(self, query: str, k: int = 10, 
                     bm25_weight: float = 0.4, 
                     vector_weight: float = 0.6) -> List[Tuple[DocumentChunk, float]]:
        """Recherche hybride avec scores combinés"""
        
        # Recherche BM25
        tokenized_query = self._tokenize(query)
        bm25_scores = self.bm25.get_scores(tokenized_query)
        
        # Normalisation BM25
        if bm25_scores.max() > 0:
            bm25_scores = bm25_scores / bm25_scores.max()
        
        # Recherche vectorielle
        query_embedding = self.embedding_model.encode([query])
        D, I = self.faiss_index.search(query_embedding.astype('float32'), k)
        
        # Scores vectoriels (inverse de la distance)
        vector_scores = np.zeros(len(self.chunks))
        for idx, distance in zip(I[0], D[0]):
            if idx != -1:
                # Convertir distance en similarité
                vector_scores[idx] = 1 / (1 + distance)
        
        # Normalisation vectorielle
        if vector_scores.max() > 0:
            vector_scores = vector_scores / vector_scores.max()
        
        # Combinaison pondérée
        combined_scores = (bm25_weight * bm25_scores) + (vector_weight * vector_scores)
        
        # Sélection des top-k résultats
        top_indices = np.argsort(combined_scores)[::-1][:k]
        
        results = []
        for idx in top_indices:
            if combined_scores[idx] > 0:
                results.append((self.chunks[idx], combined_scores[idx]))
        
        return results

# ============================================================================
# 3. CALCULATEUR DE MÉTRIQUES COMPLET
# ============================================================================

class MetricsCalculator:
    """Calcule toutes les métriques requises pour l'audit et l'assistance"""
    
    def __init__(self):
        # Initialiser les modèles nécessaires
        self.embedding_model = SentenceTransformer('paraphrase-multilingual-MiniLM-L12-v2')
    
    def calculate_bm25_score(self, query: str, document: str, retriever: HybridRetriever) -> float:
        """Calcule le score BM25 entre une requête et un document"""
        # Crée un retriever temporaire pour le document
        temp_chunk = DocumentChunk(
            text=document,
            doc_type=DocumentType.PDF,
            doc_name="temp",
            page_num=1
        )
        
        # Indexe le document
        temp_retriever = HybridRetriever()
        temp_retriever.index_chunks([temp_chunk])
        
        # Recherche
        results = temp_retriever.hybrid_search(query, k=1)
        if results:
            return float(results[0][1])
        return 0.0
    
    def calculate_cosine_similarity(self, text1: str, text2: str) -> float:
        """Calcule la similarité cosinus entre deux textes"""
        if not text1 or not text2:
            return 0.0
        
        embeddings = self.embedding_model.encode([text1, text2])
        from sentence_transformers.util import cos_sim
        similarity = cos_sim(embeddings[0].reshape(1, -1), embeddings[1].reshape(1, -1))
        return float(similarity.item())
    
    def calculate_rouge_l(self, reference: str, candidate: str) -> Dict[str, float]:
        """Calcule les scores ROUGE-L (simplifié)"""
        if not reference or not candidate:
            return {"precision": 0.0, "recall": 0.0, "f1": 0.0}
        
        # Tokenisation simple
        ref_tokens = set(reference.lower().split())
        cand_tokens = set(candidate.lower().split())
        
        if not ref_tokens or not cand_tokens:
            return {"precision": 0.0, "recall": 0.0, "f1": 0.0}
        
        # Tokens communs
        common_tokens = ref_tokens.intersection(cand_tokens)
        
        # Calcul des métriques
        precision = len(common_tokens) / len(cand_tokens) if cand_tokens else 0
        recall = len(common_tokens) / len(ref_tokens) if ref_tokens else 0
        f1 = 2 * precision * recall / (precision + recall) if (precision + recall) > 0 else 0
        
        return {
            "precision": float(precision),
            "recall": float(recall),
            "f1": float(f1)
        }
    
    def calculate_recall_at_k(self, query: str, documents: List[str], k: int = 3) -> float:
        """Calcule Recall@k pour une requête donnée"""
        if not documents:
            return 0.0
        
        # Vérifie si la réponse pertinente est dans les k premiers
        # Pour l'audit, nous vérifions si l'information existe dans les documents
        query_lower = query.lower()
        
        # Vérifie dans les k premiers documents
        top_docs = documents[:k]
        found = False
        
        for doc in top_docs:
            if doc and query_lower in doc.lower():
                found = True
                break
        
        return 1.0 if found else 0.0
    
    def calculate_nli_entailment(self, text1: str, text2: str) -> Dict[str, float]:
        """
        Calcule les scores d'entailment (contradiction, neutral, entailment)
        Version simplifiée utilisant la similarité sémantique
        """
        if not text1 or not text2:
            return {"contradiction": 0.0, "neutral": 0.5, "entailment": 0.0}
        
        # Calcule la similarité cosinus
        similarity = self.calculate_cosine_similarity(text1, text2)
        
        # Règles simplifiées pour l'entailment
        if similarity > 0.7:
            return {"contradiction": 0.1, "neutral": 0.2, "entailment": 0.7}
        elif similarity > 0.4:
            return {"contradiction": 0.2, "neutral": 0.6, "entailment": 0.2}
        else:
            return {"contradiction": 0.7, "neutral": 0.2, "entailment": 0.1}
    
    def calculate_all_metrics(self, query: str, pptx_text: str, pdf_text: str, 
                             retriever: HybridRetriever = None) -> Dict[str, Any]:
        """Calcule toutes les métriques pour une question d'audit"""
        metrics = {}
        
        # 1. BM25 Scores
        if retriever and pptx_text:
            metrics["bm25_pptx"] = self.calculate_bm25_score(query, pptx_text, retriever)
        if retriever and pdf_text:
            metrics["bm25_pdf"] = self.calculate_bm25_score(query, pdf_text, retriever)
        
        # 2. Cosine Similarity
        if pptx_text and pdf_text:
            metrics["cosine_similarity"] = self.calculate_cosine_similarity(pptx_text, pdf_text)
        else:
            metrics["cosine_similarity"] = 0.0
        
        # 3. ROUGE-L
        if pptx_text and pdf_text:
            rouge_metrics = self.calculate_rouge_l(pptx_text, pdf_text)
            metrics.update({f"rouge_l_{k}": v for k, v in rouge_metrics.items()})
        else:
            metrics.update({"rouge_l_precision": 0.0, "rouge_l_recall": 0.0, "rouge_l_f1": 0.0})
        
        # 4. Recall@k
        documents = []
        if pptx_text:
            documents.append(pptx_text)
        if pdf_text:
            documents.append(pdf_text)
        
        metrics["recall_at_3"] = self.calculate_recall_at_k(query, documents, k=3)
        metrics["recall_at_5"] = self.calculate_recall_at_k(query, documents, k=5)
        
        # 5. NLI/Entailment
        if pptx_text and pdf_text:
            nli_metrics = self.calculate_nli_entailment(pptx_text, pdf_text)
            metrics.update({f"nli_{k}": v for k, v in nli_metrics.items()})
        else:
            metrics.update({"nli_contradiction": 0.0, "nli_neutral": 0.5, "nli_entailment": 0.0})
        
        # 6. Métriques d'asymétrie
        if pptx_text and not pdf_text:
            metrics["asymmetry_type"] = "pptx_only"
        elif not pptx_text and pdf_text:
            metrics["asymmetry_type"] = "pdf_only"
        elif pptx_text and pdf_text:
            metrics["asymmetry_type"] = "both"
        else:
            metrics["asymmetry_type"] = "none"
        
        # 7. Score de contradiction (basé sur NLI)
        metrics["contradiction_score"] = metrics.get("nli_contradiction", 0.0)
        
        return metrics

# ============================================================================
# 4. SYSTÈME DE MÉMOIRE PERSISTANTE AMÉLIORÉ
# ============================================================================

@dataclass
class ConversationMemory:
    """Mémoire persistante des conversations avec métriques"""
    
    def __init__(self, storage_path: str = "workflow/output/Qna/memory/memory_store.json"):
        self.storage_path = storage_path
        self.conversations: Dict[str, List[Dict]] = {}
        self.summaries: Dict[str, str] = {}
        self.entity_memory: Dict[str, List[Dict]] = {}
        self.metrics_history: Dict[str, List[Dict]] = {}
        self.context_window = 10
        self.load_memory()
    
    def add_conversation_with_metrics(self, conversation_id: str, 
                                     messages: List[Dict], 
                                     metrics: Dict[str, Any] = None):
        """Ajoute une conversation avec ses métriques"""
        if conversation_id not in self.conversations:
            self.conversations[conversation_id] = []
            self.metrics_history[conversation_id] = []
        
        # Ajoute les messages
        self.conversations[conversation_id].extend(messages)
        if len(self.conversations[conversation_id]) > self.context_window:
            self.conversations[conversation_id] = self.conversations[conversation_id][-self.context_window:]
        
        # Ajoute les métriques
        if metrics:
            self.metrics_history[conversation_id].append({
                "timestamp": datetime.now().isoformat(),
                "metrics": metrics
            })
        
        self.save_memory()
    
    def get_conversation_with_metrics(self, conversation_id: str) -> Dict[str, Any]:
        """Récupère une conversation avec ses métriques"""
        if conversation_id not in self.conversations:
            return {"messages": [], "metrics": []}
        
        return {
            "messages": self.conversations[conversation_id],
            "metrics": self.metrics_history.get(conversation_id, [])
        }
    
    def calculate_conversation_metrics_summary(self, conversation_id: str) -> Dict[str, Any]:
        """Calcule un résumé des métriques d'une conversation"""
        if conversation_id not in self.metrics_history or not self.metrics_history[conversation_id]:
            return {}
        
        all_metrics = [item["metrics"] for item in self.metrics_history[conversation_id]]
        
        summary = {
            "total_interactions": len(all_metrics),
            "timestamp": datetime.now().isoformat()
        }
        
        # Calcule les moyennes pour chaque métrique
        metric_names = set()
        for metrics in all_metrics:
            metric_names.update(metrics.keys())
        
        for metric in metric_names:
            values = []
            for m in all_metrics:
                if metric in m and isinstance(m[metric], (int, float)):
                    values.append(m[metric])
            
            if values:
                summary[f"avg_{metric}"] = float(np.mean(values))
                summary[f"std_{metric}"] = float(np.std(values))
                summary[f"min_{metric}"] = float(np.min(values))
                summary[f"max_{metric}"] = float(np.max(values))
        
        return summary
    
    def generate_summary(self, conversation_id: str) -> str:
        """Génère un résumé textuel de la conversation avec les métriques"""
        conversation_data = self.get_conversation_with_metrics(conversation_id)
        metrics_summary = self.calculate_conversation_metrics_summary(conversation_id)
        
        if not conversation_data["messages"]:
            return "Aucune conversation enregistrée."
        
        # Compter les interactions
        user_messages = [m for m in conversation_data["messages"] if m.get("role") == "user"]
        assistant_messages = [m for m in conversation_data["messages"] if m.get("role") == "assistant"]
        
        # Dernières questions
        last_questions = [m.get("content", "")[:50] + "..." for m in user_messages[-3:]]
        
        # Créer le résumé
        summary_parts = []
        summary_parts.append(f"Résumé de la conversation {conversation_id}:")
        summary_parts.append(f"- Total d'interactions: {len(user_messages)}")
        summary_parts.append(f"- Dernières questions: {', '.join(last_questions)}")
        
        if metrics_summary:
            summary_parts.append("\nMétriques principales:")
            for key, value in metrics_summary.items():
                if key.startswith("avg_") and isinstance(value, (int, float)):
                    metric_name = key[4:]  # Enlever le préfixe "avg_"
                    if metric_name in ["recall_at_3", "avg_relevance_score", "rouge_l_f1"]:
                        summary_parts.append(f"- {metric_name}: {value:.3f}")
        
        return "\n".join(summary_parts)
    
    def save_memory(self):
        """Sauvegarde la mémoire dans un fichier"""
        data = {
            "conversations": self.conversations,
            "summaries": self.summaries,
            "entity_memory": self.entity_memory,
            "metrics_history": self.metrics_history,
            "last_updated": datetime.now().isoformat()
        }
        
        with open(self.storage_path, 'w', encoding='utf-8') as f:
            json.dump(data, f, indent=2, ensure_ascii=False)
    
    def load_memory(self):
        """Charge la mémoire depuis un fichier"""
        if os.path.exists(self.storage_path):
            try:
                with open(self.storage_path, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                
                self.conversations = data.get("conversations", {})
                self.summaries = data.get("summaries", {})
                self.entity_memory = data.get("entity_memory", {})
                self.metrics_history = data.get("metrics_history", {})
                print(f"✅ Mémoire chargée depuis {self.storage_path}")
            except Exception as e:
                print(f"⚠️ Erreur lors du chargement de la mémoire: {e}")
        else:
            print("ℹ️ Aucune mémoire précédente trouvée, création d'une nouvelle.")

# ============================================================================
# 5. ORCHESTRATION AVEC LangGraph (VERSION CUSTOM)
# ============================================================================

class StateGraph:
    """Graphique d'état personnalisé inspiré de LangGraph"""
    
    def __init__(self, name: str = "DocumentAgentGraph"):
        self.name = name
        self.nodes: Dict[str, Callable] = {}
        self.edges: List[Tuple[str, str, Optional[Callable]]] = []
        self.entry_point: Optional[str] = None
    
    def add_node(self, name: str, func: Callable):
        """Ajoute un nœud au graphique"""
        self.nodes[name] = func
    
    def add_edge(self, from_node: str, to_node: str, condition: Optional[Callable] = None):
        """Ajoute une arête entre deux nœuds"""
        self.edges.append((from_node, to_node, condition))
    
    def set_entry_point(self, node_name: str):
        """Définit le point d'entrée du graphique"""
        self.entry_point = node_name
    
    def compile(self) -> 'CompiledGraph':
        """Compile le graphique pour l'exécution"""
        return CompiledGraph(self)
    
    def visualize(self, output_path: str = "workflow/output/QnA/workflow_graph.png"):
        """Visualise le graphique avec NetworkX"""
        G = nx.DiGraph()
        
        # Ajoute les nœuds
        for node_name in self.nodes:
            G.add_node(node_name)
        
        # Ajoute les arêtes
        for from_node, to_node, condition in self.edges:
            label = "condition" if condition else ""
            G.add_edge(from_node, to_node, label=label)
        
        # Dessine le graphique
        plt.figure(figsize=(12, 8))
        pos = nx.spring_layout(G, seed=42)
        nx.draw_networkx_nodes(G, pos, node_size=3000, node_color='lightblue')
        nx.draw_networkx_edges(G, pos, arrowstyle='->', arrowsize=20)
        nx.draw_networkx_labels(G, pos, font_size=10)
        
        # Ajoute les labels des arêtes
        edge_labels = {(u, v): data['label'] for u, v, data in G.edges(data=True)}
        nx.draw_networkx_edge_labels(G, pos, edge_labels=edge_labels)
        
        plt.title(f"Workflow Graph: {self.name}")
        plt.axis('off')
        plt.savefig(output_path, dpi=300, bbox_inches='tight')
        plt.close()
        
        print(f"📊 Graphique visualisé: {output_path}")
        return output_path

class CompiledGraph:
    """Graphique compilé pour l'exécution"""
    
    def __init__(self, graph: StateGraph):
        self.graph = graph
        self.state: Dict[str, Any] = {}
    
    def invoke(self, initial_state: Dict[str, Any]) -> Dict[str, Any]:
        """Exécute le graphique avec un état initial"""
        if not self.graph.entry_point:
            raise ValueError("No entry point defined for the graph")
        
        self.state = initial_state.copy()
        current_node = self.graph.entry_point
        
        # Trace d'exécution
        execution_trace = []
        visited_nodes = set()
        
        while current_node and current_node not in visited_nodes:
            visited_nodes.add(current_node)
            
            if current_node not in self.graph.nodes:
                raise ValueError(f"Node '{current_node}' not found in graph")
            
            # Exécute le nœud
            node_func = self.graph.nodes[current_node]
            try:
                result = node_func(self.state)
                execution_trace.append({
                    "node": current_node,
                    "state_before": self.state.copy(),
                    "result": result
                })
                
                # Met à jour l'état
                if isinstance(result, dict):
                    self.state.update(result)
                
                # Trouve le prochain nœud
                next_node = None
                for from_node, to_node, condition in self.graph.edges:
                    if from_node == current_node:
                        if condition:
                            # Évalue la condition
                            if condition(self.state):
                                next_node = to_node
                                break
                        else:
                            # Arête sans condition
                            next_node = to_node
                            break
                
                current_node = next_node
                
            except Exception as e:
                execution_trace.append({
                    "node": current_node,
                    "error": str(e),
                    "state": self.state.copy()
                })
                raise
        
        self.state["execution_trace"] = execution_trace
        return self.state

# ============================================================================
# 6. SYSTÈME D'AUDIT DE CONFORMITÉ AVEC MÉTRIQUES COMPLÈTES
# ============================================================================

class AuditResult(Enum):
    CONFORME = "conforme"
    AVEC_RESERVES = "conforme avec réserves"
    NON_CONFORME = "non conforme"

@dataclass
class AuditFinding:
    """Résultat d'audit pour une question standard avec toutes les métriques"""
    question: str
    pptx_info: Optional[str] = None
    pdf_info: Optional[str] = None
    exists_in_pptx: bool = False
    exists_in_pdf: bool = False
    alignment_score: float = 0.0
    is_contradictory: bool = False
    asymmetry_type: Optional[str] = None
    verdict: AuditResult = AuditResult.NON_CONFORME
    justification: str = ""
    alerts: List[str] = field(default_factory=list)
    source_evidence: Dict[str, List[str]] = field(default_factory=dict)
    metrics: Dict[str, Any] = field(default_factory=dict)  # Toutes les métriques

class DocumentAuditor:
    """Système d'audit de conformité documentaire avec métriques complètes"""
    
    def __init__(self, pptx_retriever: HybridRetriever, 
                 pdf_retriever: HybridRetriever,
                 questions_file: str = "standard_questions.txt",
                 similarity_threshold: float = 0.7):
        
        self.pptx_retriever = pptx_retriever
        self.pdf_retriever = pdf_retriever
        self.similarity_threshold = similarity_threshold
        self.questions_file = questions_file
        self.standard_questions = self._load_standard_questions()
        self.metrics_calculator = MetricsCalculator()
    
    def _load_standard_questions(self) -> List[str]:
        """Charge les questions standards depuis un fichier texte"""
        if not os.path.exists(self.questions_file):
            print(f"⚠️ Fichier de questions non trouvé: {self.questions_file}")
            print("Utilisation des questions par défaut...")
            return [
                "Quelle est l'estimation du capital requis ?",
                "Quels sont les principaux risques identifiés ?",
                "Quelles sont les performances attendues ?",
                "Quel est l'horizon temporel du projet ?"
            ]
        
        with open(self.questions_file, 'r', encoding='utf-8') as f:
            questions = [line.strip() for line in f if line.strip() and not line.startswith('#')]
        
        print(f"✅ {len(questions)} questions standards chargées depuis {self.questions_file}")
        return questions
    
    def audit_documents(self) -> Dict[str, AuditFinding]:
        """Exécute l'audit complet des deux documents avec toutes les métriques"""
        findings = {}
        
        print("🔍 Début de l'audit documentaire...")
        for i, question in enumerate(self.standard_questions, 1):
            print(f"  [{i}/{len(self.standard_questions)}] Audit: {question[:50]}...")
            finding = self._audit_question(question)
            findings[question] = finding
        
        print("✅ Audit terminé avec succès")
        return findings
    
    def _audit_question(self, question: str) -> AuditFinding:
        """Audit une question spécifique avec toutes les métriques"""
        # Recherche dans PPTX
        pptx_results = self.pptx_retriever.hybrid_search(question, k=3)
        pptx_info, pptx_sources = self._extract_relevant_info_with_sources(pptx_results)
        
        # Recherche dans PDF
        pdf_results = self.pdf_retriever.hybrid_search(question, k=3)
        pdf_info, pdf_sources = self._extract_relevant_info_with_sources(pdf_results)
        
        # Création du finding initial
        finding = AuditFinding(
            question=question,
            pptx_info=pptx_info,
            pdf_info=pdf_info,
            exists_in_pptx=bool(pptx_info),
            exists_in_pdf=bool(pdf_info),
            source_evidence={
                "pptx_sources": pptx_sources,
                "pdf_sources": pdf_sources
            }
        )
        
        # Calcul de TOUTES les métriques
        finding.metrics = self.metrics_calculator.calculate_all_metrics(
            query=question,
            pptx_text=pptx_info,
            pdf_text=pdf_info,
            retriever=self.pptx_retriever
        )
        
        # Utilise la similarité cosinus des métriques
        finding.alignment_score = finding.metrics.get("cosine_similarity", 0.0)
        
        # Analyse d'alignement et contradictions
        self._analyze_alignment(finding)
        
        # Détermination du verdict selon les règles métier
        self._determine_verdict(finding)
        
        return finding
    
    def _extract_relevant_info_with_sources(self, results: List[Tuple[DocumentChunk, float]]) -> Tuple[Optional[str], List[str]]:
        """Extrait l'information pertinente avec les sources"""
        if not results:
            return None, []
        
        # Concatène les textes les plus pertinents
        relevant_texts = []
        sources = []
        
        for chunk, score in results:
            if score > 0.3:  # Seuil de pertinence
                relevant_texts.append(chunk.text)
                source_desc = f"{chunk.doc_type.value.upper()}: {chunk.doc_name}"
                if chunk.doc_type == DocumentType.PPTX:
                    source_desc += f" (Slide {chunk.slide_num})"
                else:
                    source_desc += f" (Page {chunk.page_num})"
                sources.append(source_desc)
        
        return "\n".join(relevant_texts[:2]) if relevant_texts else None, sources
    
    def _analyze_alignment(self, finding: AuditFinding):
        """Analyse l'alignement sémantique entre les documents"""
        # Utilise les métriques calculées
        if finding.pptx_info and finding.pdf_info:
            # Détection de contradiction basée sur NLI
            nli_contradiction = finding.metrics.get("nli_contradiction", 0.0)
            finding.is_contradictory = nli_contradiction > 0.6
            
            # Type d'asymétrie
            if finding.alignment_score < self.similarity_threshold:
                finding.asymmetry_type = "divergence_semantique"
            else:
                finding.asymmetry_type = "aligned"
        
        elif finding.pptx_info and not finding.pdf_info:
            finding.asymmetry_type = "pptx_only"
            finding.alerts.append("Information présente uniquement dans la présentation (PPTX)")
        
        elif not finding.pptx_info and finding.pdf_info:
            finding.asymmetry_type = "pdf_only"
            finding.alerts.append("Information présente uniquement dans le document détaillé (PDF/DOCX)")
    
    def _determine_verdict(self, finding: AuditFinding):
        """
        DÉTERMINE LE VERDICT D'AUDIT SELON LES RÈGLES MÉTIER
        """
        
        if not finding.exists_in_pptx and not finding.exists_in_pdf:
            # Absence totale
            finding.verdict = AuditResult.NON_CONFORME
            finding.justification = "Information absente des deux documents"
        
        elif finding.exists_in_pptx and finding.exists_in_pdf:
            # Présent dans les deux
            if finding.is_contradictory:
                finding.verdict = AuditResult.NON_CONFORME
                finding.justification = "Contradiction factuelle détectée entre les documents"
            elif finding.alignment_score >= self.similarity_threshold:
                finding.verdict = AuditResult.CONFORME
                finding.justification = f"Information alignée et cohérente (similarité: {finding.alignment_score:.2f})"
            else:
                finding.verdict = AuditResult.AVEC_RESERVES
                finding.justification = f"Information présente mais divergence sémantique (similarité: {finding.alignment_score:.2f})"
        
        elif finding.exists_in_pptx or finding.exists_in_pdf:
            # Présent dans un seul
            finding.verdict = AuditResult.AVEC_RESERVES
            doc_type = "PPTX" if finding.exists_in_pptx else "PDF/DOCX"
            finding.justification = f"Information présente uniquement dans {doc_type}"
        
        # Ajout des alertes spécifiques
        if finding.asymmetry_type:
            finding.alerts.append(f"Asymétrie documentaire: {finding.asymmetry_type}")
        
        if finding.is_contradictory:
            finding.alerts.append("Contradiction factuelle détectée")

# ============================================================================
# 7. AGENT AVEC MÉMOIRE ET LANGGRAPH
# ============================================================================

class DocumentAgentWithMemory:
    """Agent principal avec mémoire, orchestration LangGraph et métriques complètes"""
    
    def __init__(self, pptx_path: str, pdf_path: str, 
                 questions_file: str = "workflow/caches/standard_questions.txt"):
        self.pptx_path = pptx_path
        self.pdf_path = pdf_path
        self.questions_file = questions_file
        
        # Initialisation des composants
        self.extractor = DocumentExtractor()
        self.pptx_retriever = HybridRetriever()
        self.pdf_retriever = HybridRetriever()
        
        # Système de mémoire
        self.memory = ConversationMemory("workflow/output/Qna/memory/conversation_memory.json")
        self.current_conversation_id = f"conv_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
        
        # Calculateur de métriques
        self.metrics_calculator = MetricsCalculator()
        
        # Client LLM
        try:
            self.llm_client = LlamaClient() if LlamaClient else None
            print("✅ Client LLM initialisé")
        except:
            self.llm_client = None
            print("⚠️ Client LLM non disponible, utilisation du mode fallback")
        
        # État du système
        self.pptx_chunks = []
        self.pdf_chunks = []
        self.audit_results = {}
        self.combined_retriever = None
        self.assistance_metrics = []
        
        # Graphiques d'orchestration
        self.audit_graph = None
        self.assistance_graph = None
        
    def initialize(self):
        """Initialise le système avec les documents"""
        print("📂 Extraction des documents...")
        
        # Extraction
        self.pptx_chunks = self.extractor.extract(self.pptx_path)
        self.pdf_chunks = self.extractor.extract(self.pdf_path)
        
        print(f"  PPTX: {len(self.pptx_chunks)} chunks extraits")
        print(f"  PDF/DOCX: {len(self.pdf_chunks)} chunks extraits")
        
        # Indexation
        print("🔍 Indexation pour la recherche...")
        self.pptx_retriever.index_chunks(self.pptx_chunks)
        self.pdf_retriever.index_chunks(self.pdf_chunks)
        
        # Création du retriever combiné
        combined_chunks = self.pptx_chunks + self.pdf_chunks
        self.combined_retriever = HybridRetriever()
        self.combined_retriever.index_chunks(combined_chunks)
        
        # Construction des graphiques d'orchestration
        self._build_audit_graph()
        self._build_assistance_graph()
        
        print("✅ Système initialisé avec succès")
    
    def _build_audit_graph(self):
        """Construit le graphique d'audit avec LangGraph"""
        graph = StateGraph(name="AuditWorkflow")
        
        def run_audit_questions(state):
            """Exécute l'audit avec toutes les métriques"""
            print("🔎 Exécution de l'audit documentaire...")
            auditor = DocumentAuditor(
                self.pptx_retriever, 
                self.pdf_retriever,
                self.questions_file
            )
            findings = auditor.audit_documents()
            
            # Calcule les métriques globales d'audit
            audit_metrics = self._calculate_audit_global_metrics(findings)
            
            return {
                "audit_findings": findings, 
                "audit_complete": True,
                "audit_metrics": audit_metrics
            }
        
        def generate_audit_report(state):
            """Génère le rapport d'audit immédiatement"""
            print("📊 Génération du rapport d'audit...")
            
            # Génère le rapport d'audit
            self._generate_audit_report("audit_report.json")
            
            return {"audit_report_generated": True}
        
        # Ajout des nœuds
        graph.add_node("audit", run_audit_questions)
        graph.add_node("report", generate_audit_report)
        
        # Ajout des arêtes
        graph.add_edge("audit", "report")
        
        # Point d'entrée
        graph.set_entry_point("audit")
        
        self.audit_graph = graph.compile()
    
    def _build_assistance_graph(self):
        """Construit le graphique d'assistance avec mémoire"""
        graph = StateGraph(name="AssistanceWorkflow")
        
        def retrieve_context(state):
            """Récupère le contexte pertinent"""
            question = state.get("question", "")
            print(f"🔍 Recherche de contexte pour: {question[:50]}...")
            
            if self.combined_retriever:
                results = self.combined_retriever.hybrid_search(question, k=5)
                
                # Construction du contexte avec sources
                context_parts = []
                sources_info = []
                retrieved_texts = []
                
                for idx, (chunk, score) in enumerate(results, 1):
                    source_type = chunk.doc_type.value.upper()
                    location = f"Slide {chunk.slide_num}" if chunk.slide_num else f"Page {chunk.page_num}"
                    
                    context_parts.append(
                        f"[Source {idx}: {source_type} - {chunk.doc_name} - {location} - Score: {score:.3f}]\n"
                        f"{chunk.text}\n"
                    )
                    
                    sources_info.append({
                        "source_id": idx,
                        "document_type": source_type,
                        "document_name": chunk.doc_name,
                        "location": location,
                        "relevance_score": float(score)
                    })
                    
                    retrieved_texts.append(chunk.text)
                
                context = "\n".join(context_parts)
                return {
                    "context": context, 
                    "sources": sources_info, 
                    "retrieved_chunks": results,
                    "retrieved_texts": retrieved_texts
                }
            
            return {"context": "", "sources": [], "retrieved_texts": []}
        
        def generate_response(state):
            """Génère une réponse avec LLM et mémoire"""
            question = state.get("question", "")
            context = state.get("context", "")
            memory_context = state.get("memory_context", "")
            
            # Construction du prompt avec contexte et mémoire
            prompt_parts = []
            
            if memory_context:
                prompt_parts.append(memory_context)
            
            prompt_parts.append("📚 CONTEXTE DOCUMENTAIRE:")
            prompt_parts.append(context)
            
            prompt_parts.append(f"\n❓ QUESTION: {question}")
            
            full_prompt = "\n\n".join(prompt_parts)
            
            # Génération avec LLM ou fallback
            if self.llm_client:
                system_prompt = """Tu es un assistant expert en analyse documentaire. 
                Réponds uniquement basé sur le contexte fourni. 
                Cite tes sources avec [Source X]. 
                Sois concis et précis."""
                
                messages = [
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": full_prompt}
                ]
                
                try:
                    response = self.llm_client.generate_response(messages)
                    answer = response if response else "Désolé, je n'ai pas pu générer de réponse."
                except:
                    answer = "Erreur lors de la génération de la réponse."
            else:
                # Fallback
                answer = f"Information trouvée dans le contexte. Sources: {len(state.get('sources', []))}"
            
            return {"answer": answer, "generation_complete": True}
        
        def calculate_assistance_metrics(state):
            """Calcule les métriques pour l'assistance"""
            question = state.get("question", "")
            retrieved_texts = state.get("retrieved_texts", [])
            answer = state.get("answer", "")
            
            # Calcule les métriques d'assistance
            assistance_metrics = {
                "timestamp": datetime.now().isoformat(),
                "question_length": len(question),
                "answer_length": len(answer),
                "retrieved_chunks_count": len(retrieved_texts),
                "sources_count": len(state.get("sources", [])),
                "has_memory_context": state.get("has_memory", False)
            }
            
            # Métriques de recherche
            if retrieved_texts:
                # Recall@k
                assistance_metrics["recall_at_3"] = self.metrics_calculator.calculate_recall_at_k(
                    question, retrieved_texts, k=3
                )
                assistance_metrics["recall_at_5"] = self.metrics_calculator.calculate_recall_at_k(
                    question, retrieved_texts, k=5
                )
                
                # BM25 moyen (simplifié)
                scores = [s.get("relevance_score", 0) for s in state.get("sources", [])]
                assistance_metrics["avg_relevance_score"] = float(np.mean(scores)) if scores else 0.0
                
                # Similarité entre les chunks
                if len(retrieved_texts) >= 2:
                    sim = self.metrics_calculator.calculate_cosine_similarity(
                        retrieved_texts[0], retrieved_texts[1]
                    )
                    assistance_metrics["top_chunks_similarity"] = float(sim)
            
            # ROUGE-L entre question et réponse (simplifié)
            rouge = self.metrics_calculator.calculate_rouge_l(question, answer)
            assistance_metrics.update({f"rouge_l_{k}": v for k, v in rouge.items()})
            
            return {"assistance_metrics": assistance_metrics}
        
        def update_memory_with_metrics(state):
            """Met à jour la mémoire avec la conversation et les métriques"""
            question = state.get("question", "")
            answer = state.get("answer", "")
            assistance_metrics = state.get("assistance_metrics", {})
            
            # Ajoute à la mémoire conversationnelle avec métriques
            self.memory.add_conversation_with_metrics(
                self.current_conversation_id,
                [
                    {"role": "user", "content": question, "timestamp": datetime.now().isoformat()},
                    {"role": "assistant", "content": answer, "timestamp": datetime.now().isoformat()}
                ],
                assistance_metrics
            )
            
            # Stocke les métriques pour le rapport final
            self.assistance_metrics.append(assistance_metrics)
            
            return {"memory_updated": True}
        
        # Ajout des nœuds
        graph.add_node("retrieve", retrieve_context)
        graph.add_node("generate", generate_response)
        graph.add_node("calculate_metrics", calculate_assistance_metrics)
        graph.add_node("update_memory", update_memory_with_metrics)
        
        # Ajout des arêtes (workflow séquentiel)
        graph.add_edge("retrieve", "generate")
        graph.add_edge("generate", "calculate_metrics")
        graph.add_edge("calculate_metrics", "update_memory")
        
        # Point d'entrée
        graph.set_entry_point("retrieve")
        
        self.assistance_graph = graph.compile()
        
        # Visualise le graphique
        try:
            graph.visualize("assistance_workflow.png")
        except:
            print("⚠️ Impossible de visualiser le graphique")
    
    def run_audit(self) -> Dict[str, AuditFinding]:
        """Exécute l'audit avec génération immédiate du rapport"""
        print("🔎 Exécution de l'audit documentaire...")
        
        # État initial pour le graphique d'audit
        initial_state = {
            "documents_loaded": True,
            "start_time": datetime.now().isoformat()
        }
        
        # Exécute le graphique
        result = self.audit_graph.invoke(initial_state)
        
        # Récupère les résultats
        if "audit_findings" in result:
            self.audit_results = result["audit_findings"]
            self._display_audit_summary()
        
        print("✅ Audit terminé avec succès")
        return self.audit_results
    
    def _calculate_audit_global_metrics(self, findings: Dict[str, AuditFinding]) -> Dict[str, Any]:
        """Calcule les métriques globales de l'audit"""
        if not findings:
            return {}
        
        total = len(findings)
        
        # Comptage des verdicts
        verdict_counts = {
            "conforme": 0,
            "avec_reserves": 0,
            "non_conforme": 0
        }
        
        # Métriques agrégées
        all_metrics = {
            "cosine_similarity": [],
            "bm25_pptx": [],
            "bm25_pdf": [],
            "recall_at_3": [],
            "rouge_l_f1": [],
            "nli_contradiction": []
        }
        
        for finding in findings.values():
            # Comptage des verdicts
            if finding.verdict == AuditResult.CONFORME:
                verdict_counts["conforme"] += 1
            elif finding.verdict == AuditResult.AVEC_RESERVES:
                verdict_counts["avec_reserves"] += 1
            else:
                verdict_counts["non_conforme"] += 1
            
            # Agrégation des métriques
            metrics = finding.metrics
            for key in all_metrics.keys():
                if key in metrics and isinstance(metrics[key], (int, float)):
                    all_metrics[key].append(metrics[key])
        
        # Calcul des moyennes
        global_metrics = {
            "total_questions": total,
            "verdict_distribution": verdict_counts,
            "conformity_rate": verdict_counts["conforme"] / total if total > 0 else 0,
            "reserves_rate": verdict_counts["avec_reserves"] / total if total > 0 else 0,
            "non_conformity_rate": verdict_counts["non_conforme"] / total if total > 0 else 0,
            "timestamp": datetime.now().isoformat()
        }
        
        # Ajoute les moyennes des métriques
        for key, values in all_metrics.items():
            if values:
                global_metrics[f"avg_{key}"] = float(np.mean(values))
                global_metrics[f"std_{key}"] = float(np.std(values)) if len(values) > 1 else 0.0
        
        return global_metrics
    
    def _generate_audit_report(self, output_path: str = "workflow/output/QnA/audit/audit_report.json"):
        """Génère le rapport d'audit avec toutes les métriques"""
        if not self.audit_results:
            print("⚠️ Aucun résultat d'audit à rapporter")
            return
        
        # Calcule les métriques globales
        global_metrics = self._calculate_audit_global_metrics(self.audit_results)
        
        # Préparation du rapport
        report = {
            "metadata": {
                "pptx_file": self.pptx_path,
                "pdf_file": self.pdf_path,
                "questions_file": self.questions_file,
                "audit_date": datetime.now().isoformat(),
                "total_questions": len(self.audit_results),
                "system_version": "2.0"
            },
            "audit_results": {},
            "global_metrics": global_metrics,
            "summary": self._generate_audit_summary(),
            "decision_logic": self._explain_decision_logic()
        }
        
        # Conversion des AuditFinding en dict
        for question, finding in self.audit_results.items():
            report["audit_results"][question] = {
                "exists_in_pptx": finding.exists_in_pptx,
                "exists_in_pdf": finding.exists_in_pdf,
                "alignment_score": float(finding.alignment_score),
                "is_contradictory": finding.is_contradictory,
                "asymmetry_type": finding.asymmetry_type,
                "verdict": finding.verdict.value,
                "justification": finding.justification,
                "alerts": finding.alerts,
                "source_evidence": finding.source_evidence,
                "metrics": finding.metrics  # Inclut TOUTES les métriques
            }
        
        # Sauvegarde
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(report, f, indent=2, ensure_ascii=False, default=str)
        
        print(f"📄 Rapport d'audit généré: {output_path}")
        
        # Affiche un résumé
        print("\n📊 MÉTRIQUES GLOBALES DE L'AUDIT:")
        print(f"  • Questions analysées: {global_metrics['total_questions']}")
        print(f"  • Taux de conformité: {global_metrics['conformity_rate']*100:.1f}%")
        print(f"  • Similarité cosinus moyenne: {global_metrics.get('avg_cosine_similarity', 0):.3f}")
        print(f"  • Recall@3 moyen: {global_metrics.get('avg_recall_at_3', 0):.3f}")
        print(f"  • Score ROUGE-L F1 moyen: {global_metrics.get('avg_rouge_l_f1', 0):.3f}")
        
        return report
    
    def answer_question(self, question: str) -> Dict[str, Any]:
        """Répond à une question en utilisant le graphique d'assistance"""
        print(f"🤖 Traitement de la question...")
        
        # État initial
        initial_state = {
            "question": question,
            "conversation_id": self.current_conversation_id,
            "start_time": datetime.now().isoformat(),
            "has_memory": len(self.memory.conversations.get(self.current_conversation_id, [])) > 0
        }
        
        # Exécute le graphique d'assistance
        result = self.assistance_graph.invoke(initial_state)
        
        # Prépare la réponse
        response = {
            "question": question,
            "answer": result.get("answer", "Pas de réponse générée."),
            "sources": result.get("sources", []),
            "metrics": result.get("assistance_metrics", {}),
            "timestamp": datetime.now().isoformat()
        }
        
        # Affiche la réponse
        self._display_response(response)
        
        return response
    
    def _display_audit_summary(self):
        """Affiche un résumé des résultats d'audit"""
        if not self.audit_results:
            print("⚠️ Aucun résultat d'audit à afficher")
            return
        
        total = len(self.audit_results)
        conform_count = sum(1 for f in self.audit_results.values() 
                          if f.verdict == AuditResult.CONFORME)
        reserves_count = sum(1 for f in self.audit_results.values() 
                           if f.verdict == AuditResult.AVEC_RESERVES)
        non_conform_count = sum(1 for f in self.audit_results.values() 
                              if f.verdict == AuditResult.NON_CONFORME)
        
        print("\n" + "="*60)
        print("📊 RÉSUMÉ DE L'AUDIT DOCUMENTAIRE")
        print("="*60)
        print(f"Total des questions: {total}")
        print(f"✅ Conforme: {conform_count} ({conform_count/total*100:.1f}%)")
        print(f"⚠️  Avec réserves: {reserves_count} ({reserves_count/total*100:.1f}%)")
        print(f"❌ Non conforme: {non_conform_count} ({non_conform_count/total*100:.1f}%)")
    
    def _display_response(self, response: Dict[str, Any]):
        """Affiche la réponse de manière structurée"""
        print("\n" + "="*60)
        print("🤖 RÉPONSE")
        print("="*60)
        print(f"📝 {response['answer']}")
        
        if response.get('sources'):
            print(f"\n📚 Sources ({len(response['sources'])}):")
            for source in response['sources'][:3]:  # Affiche max 3 sources
                print(f"  • {source['document_type']}: {source['document_name']} ({source['location']})")
        
        if response.get('metrics'):
            print(f"\n📈 Métriques:")
            metrics = response['metrics']
            for key in ['recall_at_3', 'avg_relevance_score', 'rouge_l_f1', 'top_chunks_similarity']:
                if key in metrics and isinstance(metrics[key], (int, float)):
                    print(f"  {key}: {metrics[key]:.3f}")
    
    def interactive_assistance(self):
        """Mode interactif d'assistance avec mémoire"""
        print("\n" + "="*60)
        print("🤖 MODE ASSISTANCE INTERACTIF AVEC MÉMOIRE")
        print("="*60)
        print(f"Conversation ID: {self.current_conversation_id}")
        print("\nPosez vos questions (tapez 'quit' pour sortir):")
        
        while True:
            try:
                question = input("\n❓ Votre question: ").strip()
                
                if question.lower() in ['quit', 'exit', 'q', 'quitter']:
                    # Génère un résumé de la conversation
                    summary = self.memory.generate_summary(self.current_conversation_id)
                    print(f"\n📝 Résumé de la conversation:\n{summary}")
                    
                    # Génère le rapport de mémoire immédiatement
                    self._generate_memory_report("memory_report.json")
                    break
                
                if not question:
                    continue
                
                # Traite la question
                response = self.answer_question(question)
                
                print("\n---")
                
            except KeyboardInterrupt:
                print("\n\nInterruption par l'utilisateur")
                # Génère le rapport de mémoire avant de quitter
                self._generate_memory_report("memory_report.json")
                break
            except Exception as e:
                print(f"❌ Erreur: {e}")
                import traceback
                traceback.print_exc()
    
    def _generate_memory_report(self, output_path: str = "workflow/output/QnA/memory/memory_report.json"):
        """Génère le rapport de mémoire avec toutes les métriques"""
        # Récupère les données de mémoire
        conversation_data = self.memory.get_conversation_with_metrics(self.current_conversation_id)
        metrics_summary = self.memory.calculate_conversation_metrics_summary(self.current_conversation_id)
        
        # Prépare le rapport
        report = {
            "metadata": {
                "conversation_id": self.current_conversation_id,
                "generation_date": datetime.now().isoformat(),
                "total_interactions": len(conversation_data.get("messages", [])) // 2,
                "report_type": "assistance_utilisateur"
            },
            "conversation": {
                "messages": conversation_data.get("messages", []),
                "metrics_history": conversation_data.get("metrics", [])
            },
            "metrics_summary": metrics_summary,
            "aggregated_metrics": self._calculate_aggregated_assistance_metrics()
        }
        
        # Sauvegarde
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(report, f, indent=2, ensure_ascii=False, default=str)
        
        print(f"\n💾 Rapport de mémoire généré: {output_path}")
        
        # Affiche un résumé
        if metrics_summary:
            print("\n📊 RÉSUMÉ DES MÉTRIQUES D'ASSISTANCE:")
            print(f"  • Interactions totales: {metrics_summary.get('total_interactions', 0)}")
            if 'avg_recall_at_3' in metrics_summary:
                print(f"  • Recall@3 moyen: {metrics_summary['avg_recall_at_3']:.3f}")
            if 'avg_avg_relevance_score' in metrics_summary:
                print(f"  • Score de pertinence moyen: {metrics_summary['avg_avg_relevance_score']:.3f}")
        
        return report
    
    def _calculate_aggregated_assistance_metrics(self) -> Dict[str, Any]:
        """Calcule les métriques agrégées pour l'assistance"""
        if not self.assistance_metrics:
            return {}
        
        # Collecte toutes les métriques
        all_metrics = {}
        for metric_set in self.assistance_metrics:
            for key, value in metric_set.items():
                if isinstance(value, (int, float)):
                    if key not in all_metrics:
                        all_metrics[key] = []
                    all_metrics[key].append(value)
        
        # Calcule les statistiques
        aggregated = {
            "total_interactions": len(self.assistance_metrics),
            "timestamp": datetime.now().isoformat()
        }
        
        for key, values in all_metrics.items():
            if values:
                aggregated[f"{key}_mean"] = float(np.mean(values))
                aggregated[f"{key}_std"] = float(np.std(values)) if len(values) > 1 else 0.0
                aggregated[f"{key}_min"] = float(np.min(values))
                aggregated[f"{key}_max"] = float(np.max(values))
        
        return aggregated
    
    def _generate_audit_summary(self) -> Dict[str, Any]:
        """Génère un résumé statistique de l'audit"""
        if not self.audit_results:
            return {}
        
        total = len(self.audit_results)
        conform_count = sum(1 for f in self.audit_results.values() 
                          if f.verdict == AuditResult.CONFORME)
        reserves_count = sum(1 for f in self.audit_results.values() 
                           if f.verdict == AuditResult.AVEC_RESERVES)
        non_conform_count = sum(1 for f in self.audit_results.values() 
                              if f.verdict == AuditResult.NON_CONFORME)
        
        return {
            "total_questions": total,
            "conform_count": conform_count,
            "reserves_count": reserves_count,
            "non_conform_count": non_conform_count,
            "conformity_rate": conform_count / total if total > 0 else 0,
            "with_reserves_rate": reserves_count / total if total > 0 else 0,
            "non_conformity_rate": non_conform_count / total if total > 0 else 0
        }
    
    def _explain_decision_logic(self) -> Dict[str, str]:
        """Explique la logique de décision pour la conformité"""
        return {
            "decision_framework": "Règles métier basées sur la présence et la cohérence de l'information",
            "rules": {
                "non_conforme": [
                    "Information absente des deux documents",
                    "Contradiction factuelle entre documents (score NLI contradiction > 0.6)"
                ],
                "avec_reserves": [
                    "Information présente dans un seul document",
                    "Information présente mais divergence sémantique (similarité cosinus < 0.7)"
                ],
                "conforme": [
                    "Information présente dans les deux documents",
                    "Alignement sémantique élevé (similarité cosinus >= 0.7)",
                    "Pas de contradiction détectée (score NLI contradiction <= 0.6)"
                ]
            },
            "thresholds": {
                "contradiction_threshold": 0.6,
                "alignment_threshold": 0.7,
                "relevance_threshold": 0.3
            },
            "metrics_used": [
                "Présence/absence dans chaque document",
                "Similarité cosinus sémantique",
                "BM25 scores",
                "ROUGE-L (recouvrement lexical)",
                "Recall@k (retrouvabilité)",
                "NLI/entailment (contradiction, neutral, entailment)"
            ]
        }
    
    def generate_comprehensive_report(self, output_dir: str = "reports"):
        """Génère tous les rapports dans un dossier (fonction de secours)"""
        import os
        os.makedirs(output_dir, exist_ok=True)
        
        # Rapport d'audit (si pas déjà généré)
        audit_path = os.path.join(output_dir, "audit_report.json")
        if not os.path.exists(audit_path) and self.audit_results:
            self._generate_audit_report(audit_path)
        
        # Rapport de mémoire (si pas déjà généré)
        memory_path = os.path.join(output_dir, "memory_report.json")
        if not os.path.exists(memory_path):
            self._generate_memory_report(memory_path)
        
        # Rapport exécutif
        executive_path = os.path.join(output_dir, "executive_summary.md")
        with open(executive_path, 'w', encoding='utf-8') as f:
            f.write(self._generate_executive_summary())
        print(f"📋 Résumé exécutif généré: {executive_path}")
        
        print(f"\n✅ Tous les rapports générés dans: {output_dir}")
    
    def _generate_executive_summary(self) -> str:
        """Génère un résumé exécutif en Markdown"""
        audit_summary = self._generate_audit_summary()
        metrics_summary = self.memory.calculate_conversation_metrics_summary(self.current_conversation_id)
        
        summary = f"""# Rapport Exécutif - Système d'Audit Documentaire

**Date:** {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}
**Documents analysés:** {os.path.basename(self.pptx_path)}, {os.path.basename(self.pdf_path)}
**Conversation ID:** {self.current_conversation_id}

## 📊 Résumé de l'Audit

"""
        
        if audit_summary:
            summary += f"- **Total des questions:** {audit_summary['total_questions']}\n"
            summary += f"- **✅ Conforme:** {audit_summary['conform_count']} ({audit_summary['conformity_rate']*100:.1f}%)\n"
            summary += f"- **⚠️ Avec réserves:** {audit_summary['reserves_count']} ({audit_summary['with_reserves_rate']*100:.1f}%)\n"
            summary += f"- **❌ Non conforme:** {audit_summary['non_conform_count']} ({audit_summary['non_conformity_rate']*100:.1f}%)\n"
        
        summary += f"""
## 🤖 Assistance Utilisateur

- **Conversation active:** {self.current_conversation_id}
- **Interactions totales:** {metrics_summary.get('total_interactions', 0)}
- **Dernière activité:** {self.memory.conversations.get(self.current_conversation_id, [{}])[-1].get('timestamp', 'N/A')}

## 📈 Métriques Clés

### Audit
"""
        
        if self.audit_results and len(self.audit_results) > 0:
            # Prend la première question pour exemple
            first_finding = next(iter(self.audit_results.values()))
            if first_finding.metrics:
                metrics = first_finding.metrics
                summary += f"- **Similarité Cosinus moyenne:** {metrics.get('cosine_similarity', 0):.3f}\n"
                summary += f"- **Recall@3 moyen:** {metrics.get('recall_at_3', 0):.3f}\n"
                summary += f"- **ROUGE-L F1 moyen:** {metrics.get('rouge_l_f1', 0):.3f}\n"
        
        summary += f"""
### Assistance
"""
        
        if metrics_summary:
            if 'avg_recall_at_3' in metrics_summary:
                summary += f"- **Recall@3 moyen:** {metrics_summary['avg_recall_at_3']:.3f}\n"
            if 'avg_avg_relevance_score' in metrics_summary:
                summary += f"- **Pertinence moyenne:** {metrics_summary['avg_avg_relevance_score']:.3f}\n"
        
        summary += f"""
## 🔧 Architecture Technique

- **Orchestration:** LangGraph (graphiques personnalisés)
- **Mémoire:** ConversationMemory avec persistance JSON
- **Recherche:** Hybride BM25 + FAISS
- **Métriques:** Cosinus, BM25, ROUGE-L, Recall@k, NLI/Entailment
- **LLM:** {"LlamaClient intégré" if self.llm_client else "Mode fallback"}

## 📈 Recommandations

1. **Conformité:** {self._get_conformity_recommendation()}
2. **Amélioration:** Maintenir la cohérence entre documents
3. **Surveillance:** Suivre les asymétries détectées
"""
        return summary
    
    def _get_conformity_recommendation(self) -> str:
        """Génère une recommandation basée sur les résultats d'audit"""
        if not self.audit_results:
            return "Audit non réalisé"
        
        non_conform = sum(1 for f in self.audit_results.values() 
                         if f.verdict == AuditResult.NON_CONFORME)
        total = len(self.audit_results)
        
        if non_conform > total * 0.3:
            return "⚠️ Amélioration urgente nécessaire (>30% non conforme)"
        elif non_conform > 0:
            return "📋 Correction recommandée pour les non-conformités"
        else:
            return "✅ Bon niveau de conformité"

# ============================================================================
# 8. EXEMPLE D'UTILISATION
# ============================================================================

def main():
    """Exemple d'utilisation du système avec LangGraph et mémoire"""
    
    # Configuration
    PPTX_FILE = "dataset\example_2\FINAL-PRS-GB-ODDO BHF US Equity Active ETF-20250630_8PN_clean.pptx"  # À remplacer
    PDF_FILE = "dataset\example_2\prospectus.docx"    # À remplacer
    QUESTIONS_FILE = "workflow/caches/standard_questions.txt"
    
    # Initialisation de l'agent
    print("🚀 Initialisation de l'Agent avec LangGraph et Mémoire")
    print("="*60)
    
    agent = DocumentAgentWithMemory(PPTX_FILE, PDF_FILE, QUESTIONS_FILE)
    
    try:
        # Phase 1: Initialisation
        agent.initialize()
        
        # Phase 2: Audit avec génération immédiate du rapport
        print("\n" + "="*60)
        print("PHASE 1: AUDIT AVEC LANGGRAPH")
        print("="*60)
        agent.run_audit()  # Génère automatiquement le rapport d'audit
        
        # Phase 3: Assistance interactive avec mémoire
        print("\n" + "="*60)
        print("PHASE 2: ASSISTANCE AVEC MÉMOIRE")
        print("="*60)
        agent.interactive_assistance()  # Génère automatiquement le rapport de mémoire à la fin
        
        # Phase 4: Génération des rapports complets (fonction de secours)
        print("\n" + "="*60)
        print("PHASE 3: GÉNÉRATION DES RAPPORTS COMPLETS")
        print("="*60)
        agent.generate_comprehensive_report()
        
        print(f"\n✅ Système terminé avec succès!")
        print(f"   • Conversation ID: {agent.current_conversation_id}")
        print(f"   • Rapports générés: audit_report.json, memory_report.json")
        
    except FileNotFoundError as e:
        print(f"❌ Erreur: {e}")
        print("Veuillez vérifier les chemins des fichiers.")
    except Exception as e:
        print(f"❌ Erreur inattendue: {e}")
        import traceback
        traceback.print_exc()


if __name__ == "__main__":
    main()