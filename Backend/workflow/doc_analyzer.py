"""
DOC_ANALYZER - Version Workflow Simplifié (LangGraph)

Agent d'analyse documentaire intelligent utilisant LLM pour l'analyse.
- Parse documents (PPTX, DOCX, PDF)
- Charge règles JSON ODDO (contextuelles + structurelles)  
- Utilise LLM avec CAG (Context-Aware Generation) pour analyse
- Orchestration via LangGraph
"""

import json
import uuid
import os
import sys
import hashlib
from typing import Dict, List, Any, Optional, TypedDict
from dataclasses import dataclass
import numpy as np
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx import Document
import zipfile
import base64
from PIL import Image
import io
from langgraph.graph import StateGraph, END

# Imports conditionnels pour robustesse
try:
    import pytesseract
except ImportError:
    pytesseract = None

try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

from datetime import datetime

# Import modules centralisés
sys.path.append('.')
from config.settings import *
from llm_client import LlamaClient

print("=" * 80)
print("🔍 DOC_ANALYZER - Version Workflow Simplifié (LangGraph)")
print("=" * 80)

# ==================== CACHE MANAGER ====================

class DocAnalyzerCache:
    """Gestionnaire de cache pour DocAnalyzer"""
    
    def __init__(self):
        self.cache_dir = os.path.join(CACHE_DIR, "doc_analyzer_cache")
        os.makedirs(self.cache_dir, exist_ok=True)
        
    def get(self, doc_id: str, rules_hash: str) -> Optional[Dict]:
        """Récupère une analyse du cache"""
        cache_path = self._get_cache_path(doc_id, rules_hash)
        if os.path.exists(cache_path):
            try:
                with open(cache_path, 'r', encoding='utf-8') as f:
                    return json.load(f)
            except Exception as e:
                print(f"⚠️ Erreur lecture cache: {e}")
        return None
        
    def set(self, doc_id: str, rules_hash: str, data: Dict):
        """Sauvegarde une analyse dans le cache"""
        cache_path = self._get_cache_path(doc_id, rules_hash)
        try:
            with open(cache_path, 'w', encoding='utf-8') as f:
                json.dump(data, f, indent=2, ensure_ascii=False)
        except Exception as e:
            print(f"⚠️ Erreur écriture cache: {e}")
            
    def _get_cache_path(self, doc_id: str, rules_hash: str) -> str:
        """Génère le chemin du fichier cache"""
        filename = f"{doc_id}_{rules_hash}.json"
        return os.path.join(self.cache_dir, filename)


# ==================== GESTION DES RÈGLES JSON ====================

class ComplianceRuleManager:
    """Charge et gère les règles de conformité depuis JSON ODDO"""
    
    def __init__(self, metadata: Dict = None):
        self.metadata = metadata or {}
        self.all_rules = {}  # Toutes les règles chargées
        self.applicable_rules = {}  # Règles filtrées selon métadonnées
        self.load_oddo_rules()
    
    def load_oddo_rules(self):
        """Charge les règles depuis les fichiers JSON ODDO"""
        print("\n📋 Chargement des règles de conformité ODDO...")
        
        all_rules_list = []
        
        # Charger règles contextuelles (RC1-RC47)
        try:
            if os.path.exists(REGLES_CONTEXTUELLES_PATH):
                with open(REGLES_CONTEXTUELLES_PATH, 'r', encoding='utf-8') as f:
                    contextual = json.load(f)
                all_rules_list.extend(contextual)
                print(f"  ✅ {len(contextual)} règles contextuelles chargées")
            else:
                print(f"  ⚠️ {REGLES_CONTEXTUELLES_PATH} introuvable")
        except Exception as e:
            print(f"  ❌ Erreur chargement règles contextuelles: {e}")
        
        # Charger règles structurelles (RS1-RS11)
        try:
            if os.path.exists(REGLES_STRUCTURELLES_PATH):
                with open(REGLES_STRUCTURELLES_PATH, 'r', encoding='utf-8') as f:
                    structural = json.load(f)
                all_rules_list.extend(structural)
                print(f"  ✅ {len(structural)} règles structurelles chargées")
            else:
                print(f"  ⚠️ {REGLES_STRUCTURELLES_PATH} introuvable")
        except Exception as e:
            print(f"  ❌ Erreur chargement règles structurelles: {e}")
        
        # Indexer toutes les règles par ID
        for rule in all_rules_list:
            rule_id = rule.get('id')
            if rule_id:
                self.all_rules[rule_id] = rule
        
        # Filtrer les règles applicables selon métadonnées
        self.applicable_rules = self._filter_applicable_rules()
        
        print(f"  📊 Total règles: {len(self.all_rules)} chargées")
        print(f"  🎯 Règles applicables: {len(self.applicable_rules)} (selon métadonnées)\n")
    
    def _filter_applicable_rules(self, only_category: str = None) -> Dict[str, Dict]:
        if not self.metadata and not only_category:
            print("  ℹ️  Pas de métadonnées - toutes les règles sont applicables par défaut")
            return self.all_rules.copy()
        
        applicable = {}
        
        # Extraire métadonnées clés
        is_professional = self.metadata.get("Le client est-il un professionnel", False)
        is_new_product = self.metadata.get("Le document fait-il référence à un nouveau Produit", False)
        is_new_strategy = self.metadata.get("Le document fait-il référence à une nouvelle Stratégie", False)
        is_sicav_oddo = self.metadata.get("Est ce que le produit fait partie de la Sicav d'Oddo", False)
        
        client_type = "professionnel" if is_professional else "retail"
        
        # Logs uniquement si pertinent
        if self.metadata:
            print(f"  📝 Métadonnées:")
            print(f"     • Client: {client_type}")
            print(f"     • Nouveau produit: {is_new_product}")
            print(f"     • Nouvelle stratégie: {is_new_strategy}")
            print(f"     • SICAV ODDO: {is_sicav_oddo}")
        
        for rule_id, rule in self.all_rules.items():
            # Filtre par catégorie (Si spécifié)
            if only_category and rule.get('categorie') != only_category:
                continue

            is_applicable = True
            rule_desc = rule.get('description', '').lower()
            
            # Exceptions spécifiques
            if rule_id == "RC1":
                applicable[rule_id] = rule
                continue
            
            # Filtre par métadonnées (pour les règles contextuelles surtout)
            if "nouveau" in rule_desc or "nouvelle" in rule_desc:
                if not (is_new_product or is_new_strategy):
                    is_applicable = False
            
            # Les règles structurelles sont généralement toujours applicables
            if rule.get('categorie') == 'structurelle':
                is_applicable = True
            
            if is_applicable:
                applicable[rule_id] = rule
        
        return applicable
    
    def format_rules_for_llm(self) -> str:
        formatted = []
        for rule_id, rule in self.applicable_rules.items():
            formatted.append(
                f"• {rule_id} [{rule.get('categorie', 'N/A')}]: {rule.get('description', 'N/A')}\n"
                f"  Exigence: {rule.get('exigence', 'N/A')}"
            )
        return "\n".join(formatted)
    
    def format_metadata_for_llm(self) -> str:
        if not self.metadata:
            return "Aucune métadonnée fournie"
        
        is_professional = self.metadata.get("Le client est-il un professionnel", False)
        client_type = "Professionnel" if is_professional else "Retail (grand public)"
        
        formatted = f"""
Type de client: {client_type}
Nouveau produit: {'Oui' if self.metadata.get("Le document fait-il référence à un nouveau Produit") else 'Non'}
Nouvelle stratégie: {'Oui' if self.metadata.get("Le document fait-il référence à une nouvelle Stratégie") else 'Non'}
SICAV ODDO: {'Oui' if self.metadata.get("Est ce que le produit fait partie de la Sicav d'Oddo") else 'Non'}
Société de gestion: {self.metadata.get("Société de Gestion", "N/A")}
"""
        return formatted.strip()


# ==================== OCR EXTRACTOR ====================

class OpenSourceOCRExtractor:
    def __init__(self):
        self.available = False
        if pytesseract:
            try:
                pytesseract.get_tesseract_version()
                self.available = True
                print("  ✅ Tesseract OCR détecté et opérationnel")
            except Exception as e:
                print(f"  ⚠️ Tesseract non disponible (erreur config): {e}")
        else:
            print("  ⚠️ Module 'pytesseract' non installé")

    def extract_text_from_image(self, image_data: bytes) -> str:
        if not self.available:
            return ""
        try:
            image = Image.open(io.BytesIO(image_data))
            custom_config = OCR_CONFIG + f' -l {OCR_LANGUAGES}'
            text = pytesseract.image_to_string(image, config=custom_config)
            return text.strip()
        except Exception as e:
            print(f"  ❌ Erreur OCR image: {e}")
            return ""


# ==================== PARSER DOCUMENTAIRE ====================

class DocumentParser:
    def __init__(self):
        self.supported_formats = ['.pptx', '.docx', '.pdf']
        self.ocr_extractor = OpenSourceOCRExtractor() if ENABLE_OCR else None
        if ENABLE_OCR:
            print(f"  👁️ OCR activé (Tesseract: {'Oui' if self.ocr_extractor.available else 'Non'})")
        else:
            print("  zzz OCR désactivé (ENABLE_OCR=False)")
    
    def parse_document(self, file_path: str) -> Dict[str, Any]:
        ext = os.path.splitext(file_path)[1].lower()
        print(f"\n📄 Parsing document: {os.path.basename(file_path)}")
        print(f"  Format: {ext}")
        
        if ext == '.pptx':
            return self._parse_pptx(file_path)
        elif ext == '.docx':
            return self._parse_docx(file_path)
        elif ext == '.pdf':
            return self._parse_pdf(file_path)
        else:
            return {"error": f"Format {ext} non supporté", "type": "unknown"}
    
    def _parse_pptx(self, file_path: str) -> Dict[str, Any]:
        try:
            prs = Presentation(file_path)
            doc_data = {
                "type": "pptx",
                "document_id": f"doc_{uuid.uuid4().hex[:8]}",
                "file_name": os.path.basename(file_path),
                "total_slides": len(prs.slides),
                "slides": []
            }
            
            for idx, slide in enumerate(prs.slides):
                slide_data = {
                    "slide_number": idx + 1,
                    "layout": slide.slide_layout.name,
                    "content": [],
                    "images": []
                }
                
                for shape_idx, shape in enumerate(slide.shapes):
                    element_id = f"slide_{idx+1}_shape_{shape_idx+1}"
                    
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_data["content"].append({
                            "id": element_id,
                            "type": "text",
                            "text": shape.text.strip()
                        })
                    
                    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        if hasattr(shape, "table"):
                            table_data = self._extract_table(shape.table)
                            slide_data["content"].append({
                                "id": f"{element_id}_table",
                                "type": "table",
                                "data": table_data
                            })
                    
                    if hasattr(shape, "image"):
                        image_text = ""
                        if ENABLE_OCR and self.ocr_extractor and self.ocr_extractor.available:
                            try:
                                if hasattr(shape.image, 'blob'):
                                    image_text = self.ocr_extractor.extract_text_from_image(shape.image.blob)
                            except Exception as e:
                                print(f"    ⚠️ Erreur extraction image slide {idx+1}: {e}")
                        
                        slide_data["images"].append({
                            "id": f"{element_id}_image",
                            "type": "image",
                            "name": shape.name,
                            "ocr_text": image_text if image_text else None
                        })
                        
                        if image_text and len(image_text) > 5:
                            slide_data["content"].append({
                                "id": f"{element_id}_ocr",
                                "type": "image_text",
                                "text": f"[TEXTE IMAGE]: {image_text}"
                            })
                
                doc_data["slides"].append(slide_data)
            
            print(f"  ✅ {len(prs.slides)} slides parsés")
            return doc_data
        
        except Exception as e:
            print(f"  ❌ Erreur parsing PPTX: {e}")
            return {"error": str(e), "type": "pptx"}
    
    def _parse_docx(self, file_path: str) -> Dict[str, Any]:
        try:
            doc = Document(file_path)
            doc_data = {
                "type": "docx",
                "document_id": f"doc_{uuid.uuid4().hex[:8]}",
                "file_name": os.path.basename(file_path),
                "content": [],
                "images": []
            }
            
            for idx, para in enumerate(doc.paragraphs):
                if para.text.strip():
                    doc_data["content"].append({
                        "id": f"para_{idx+1}",
                        "type": "paragraph",
                        "text": para.text.strip(),
                        "style": para.style.name if para.style else "Normal"
                    })
            
            for idx, table in enumerate(doc.tables):
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                
                doc_data["content"].append({
                    "id": f"table_{idx+1}",
                    "type": "table",
                    "data": table_data
                })
            
            doc_data["images"] = self._extract_images_from_docx(file_path)
            
            if ENABLE_OCR and self.ocr_extractor and self.ocr_extractor.available:
                print(f"  👁️ Traitement OCR de {len(doc_data['images'])} images...")
                for img in doc_data["images"]:
                    try:
                        with zipfile.ZipFile(file_path, 'r') as docx_zip:
                            with docx_zip.open(img["path_in_archive"]) as img_file:
                                img_data = img_file.read()
                                text = self.ocr_extractor.extract_text_from_image(img_data)
                                img["ocr_text"] = text
                                
                                if text and len(text) > 5:
                                    doc_data["content"].append({
                                        "type": "image_text",
                                        "text": f"[TEXTE IMAGE {img['file_name']}]: {text}"
                                    })
                    except Exception as e:
                        print(f"    ⚠️ Erreur OCR image {img['file_name']}: {e}")
            
            print(f"  ✅ {len(doc_data['content'])} éléments parsés")
            return doc_data
        
        except Exception as e:
            print(f"  ❌ Erreur parsing DOCX: {e}")
            return {"error": str(e), "type": "docx"}
    
    def _parse_pdf(self, file_path: str) -> Dict[str, Any]:
        if not PyPDF2:
            return {"error": "Module PyPDF2 manquant", "type": "pdf"}
            
        try:
            doc_data = {
                "type": "pdf",
                "document_id": f"doc_{uuid.uuid4().hex[:8]}",
                "file_name": os.path.basename(file_path),
                "content": [],
                "metadata": {}
            }
            
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                if reader.metadata:
                    doc_data["metadata"] = {k: v for k, v in reader.metadata.items()}
                
                doc_data["total_pages"] = len(reader.pages)
                
                for i, page in enumerate(reader.pages):
                    text = page.extract_text()
                    if text.strip():
                        doc_data["content"].append({
                            "id": f"page_{i+1}",
                            "type": "page_text",
                            "page_number": i + 1,
                            "text": text.strip()
                        })
            
            print(f"  ✅ {doc_data['total_pages']} pages PDF parsées")
            return doc_data
            
        except Exception as e:
            print(f"  ❌ Erreur parsing PDF: {e}")
            return {"error": str(e), "type": "pdf"}
    
    def _extract_images_from_docx(self, docx_path: str) -> List[Dict]:
        images = []
        try:
            with zipfile.ZipFile(docx_path, 'r') as docx_zip:
                media_files = [f for f in docx_zip.namelist() if f.startswith('word/media/')]
                for media_file in media_files:
                    if media_file.lower().endswith(('.png', '.jpg', '.jpeg', '.gif')):
                        image_info = {
                            "file_name": os.path.basename(media_file),
                            "format": os.path.splitext(media_file)[1].lower(),
                            "path_in_archive": media_file
                        }
                        images.append(image_info)
        except Exception as e:
            print(f"  ⚠️ Erreur extraction images DOCX: {e}")
        return images
    
    def _extract_table(self, table) -> List[List[str]]:
        table_data = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                row_data.append(cell.text.strip() if hasattr(cell, 'text') else "")
            table_data.append(row_data)
        return table_data


# ==================== STATE LANGGRAPH ====================

class DocAnalyzerState(TypedDict):
    """État du graphe DocAnalyzer"""
    file_path: str
    metadata_path: Optional[str]
    metadata: Optional[Dict]  # Direct metadata dict from frontend
    
    # Intermédiaire
    doc_structure: Dict
    rules_hash: str
    doc_id: str
    
    # Résultats
    analysis: Dict
    metrics: Dict
    final_result: Dict

# ==================== AGENT PRINCIPAL (LANGGRAPH) ====================

class DocumentComplianceAgent:
    """Agent principal d'analyse de conformité documentaire avec LangGraph"""
    
    def __init__(self, metadata_path: str = None):
        print("\n" + "=" * 80)
        print("🚀 INITIALISATION DOC_ANALYZER AGENT (LANGGRAPH)")
        print("=" * 80)
        
        self.metadata_path = metadata_path
        self.parser = DocumentParser()
        self.llm_client = LlamaClient()
        self.cache = DocAnalyzerCache()
        
        # Initialisation RuleManager (sera mis à jour dans le graphe)
        self.rule_manager = ComplianceRuleManager()
        
        # Construction du Graphe
        self.workflow = self._build_workflow()
        
        print("\n✅ Agent initialisé avec succès\n")
    
    def _build_workflow(self) -> StateGraph:
        """Construit le workflow LangGraph"""
        workflow = StateGraph(DocAnalyzerState)
        
        # Nodes
        workflow.add_node("parse_document", self._parse_document_node)
        workflow.add_node("load_rules", self._load_rules_node)
        workflow.add_node("check_cache", self._check_cache_node)
        workflow.add_node("analyze_llm", self._analyze_llm_node)
        workflow.add_node("calculate_metrics", self._calculate_metrics_node)
        
        # Edges
        workflow.set_entry_point("parse_document")
        workflow.add_edge("parse_document", "load_rules")
        workflow.add_edge("load_rules", "check_cache")
        
        def cache_condition(state):
            if state.get("analysis"):
                return "calculate_metrics" # Cache hit -> Skip LLM
            return "analyze_llm"
            
        workflow.add_conditional_edges(
            "check_cache",
            cache_condition,
            ["calculate_metrics", "analyze_llm"]
        )
        
        workflow.add_edge("analyze_llm", "calculate_metrics")
        workflow.add_edge("calculate_metrics", END)
        
        return workflow.compile()

    # ==================== NODES ====================

    def _parse_document_node(self, state: DocAnalyzerState) -> DocAnalyzerState:
        """Parse le document"""
        print("📄 Node: Parsing Document")
        doc_structure = self.parser.parse_document(state["file_path"])
        
        if "error" in doc_structure:
            print(f"❌ Erreur parsing: {doc_structure['error']}")
            return {"doc_structure": doc_structure, "final_result": {"error": doc_structure['error']}}
            
        # Sauvegarde JSON parsé
        output_filename = f"{os.path.splitext(os.path.basename(state['file_path']))[0]}_parsed.json"
        output_path = os.path.join(OUTPUT_DIR, "parsed_docs", output_filename)
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(doc_structure, f, indent=2, ensure_ascii=False)
            
        return {"doc_structure": doc_structure}

    def _load_rules_node(self, state: DocAnalyzerState) -> DocAnalyzerState:
        """Charge les règles et métadonnées"""
        print("📋 Node: Chargement Règles")
        
        # Charger métadonnées - priorité au metadata direct, sinon depuis fichier
        metadata = state.get("metadata")
        if not metadata:
            path = state.get("metadata_path") or self.metadata_path
            if path and os.path.exists(path):
                try:
                    with open(path, 'r', encoding='utf-8') as f:
                        metadata = json.load(f)
                except Exception:
                    pass
        
        if metadata:
            print(f"  📝 Métadonnées utilisées: {metadata}")
        
        # Mise à jour RuleManager - Charger TOUTES les règles (structurelles ET contextuelles)
        self.rule_manager.metadata = metadata or {}
        self.rule_manager.applicable_rules = self.rule_manager._filter_applicable_rules()  # No filter = all rules
        
        # Hash pour cache (toutes les règles applicables)
        all_applicable = list(self.rule_manager.applicable_rules.values())
        rules_str = json.dumps(all_applicable, sort_keys=True)
        rules_hash = hashlib.md5(rules_str.encode('utf-8')).hexdigest()
        
        return {
            "rules_hash": rules_hash,
            "doc_id": state["doc_structure"].get("document_id", "unknown")
        }

    def _check_cache_node(self, state: DocAnalyzerState) -> DocAnalyzerState:
        """Vérifie le cache"""
        print("💾 Node: Vérification Cache")
        cached = self.cache.get(state["doc_id"], state["rules_hash"])
        if cached:
            print("  ✅ Résultat trouvé en cache")
            return {"analysis": cached}
        return {}

    def _analyze_llm_node(self, state: DocAnalyzerState) -> DocAnalyzerState:
        """Analyse via LLM"""
        print("🧠 Node: Analyse LLM")
        
        # Prompt
        doc_summary = self._summarize_document(state["doc_structure"])
        rules_context = self.rule_manager.format_rules_for_llm()
        metadata_context = self.rule_manager.format_metadata_for_llm()
        
        prompt = f"""
ANALYSE DE CONFORMITÉ STRUCTURELLE

=== MÉTADONNÉES DU DOCUMENT ===
{metadata_context}

=== RÈGLES DE CONFORMITÉ APPLICABLES ===
{rules_context}

=== CONTENU DU DOCUMENT (Avec Identifiants) ===
Type: {state['doc_structure'].get('type', 'unknown')}
ID: {state['doc_structure'].get('document_id', 'N/A')}

{doc_summary}

=== TÂCHE ===
Analyse ce document et identifie:

1. **Éléments conformes**: Quels éléments respectent les règles
2. **Éléments non conformes**: Quels éléments violent les règles et pourquoi
3. **Éléments manquants**: Quels éléments requis sont absents
4. **Score de conformité global**: 0-100%

IMPORTANT:
- Utilise les métadonnées pour adapter ton analyse.
- **CITE TOUJOURS L'IDENTIFIANT (id/page/slide)** de l'élément concerné.
- Pour PDF: cite "Page X".
- Pour PPTX: cite "Slide X - ID Y".
- Pour DOCX: cite "ID Y".

Réponds UNIQUEMENT en JSON avec cette structure:
{{
  "conformite_globale": {{
    "score": 85,
    "niveau": "bon|moyen|faible"
  }},
  "elements_conformes": [
    {{"element": "Police du titre", "location": "Slide 1 - slide_1_shape_1", "regle_id": "RS1", "justification": "..."}}
  ],
  "elements_non_conformes": [
    {{"element": "Mise en page", "location": "Page 3", "regle_id": "RS2", "violation": "...", "gravite": "haute"}}
  ],
  "elements_manquants": [
    {{"element_requis": "...", "regle_id": "RS3"}}
  ],
  "recommandations": ["...", "..."]
}}
"""
        messages = [
            {"role": "system", "content": SYSTEM_PROMPT_DOCUMENT_ANALYSIS},
            {"role": "user", "content": prompt}
        ]
        
        response = self.llm_client.generate_response(
            messages=messages,
            temperature=LLM_TEMPERATURE,
            max_tokens=LLM_MAX_TOKENS
        )
        
        analysis = self._parse_llm_response(response)
        
        # Sauvegarde Cache
        self.cache.set(state["doc_id"], state["rules_hash"], analysis)
        
        return {"analysis": analysis}

    def _calculate_metrics_node(self, state: DocAnalyzerState) -> DocAnalyzerState:
        """Calcule les métriques finales"""
        print("📊 Node: Calcul Métriques")
        analysis = state["analysis"]
        
        # Calcul métriques
        metrics = ComplianceMetrics.calculate_metrics(analysis)
        
        # Résultat final
        final_result = {
            "document_id": state["doc_id"],
            "file_name": os.path.basename(state["file_path"]),
            "timestamp": datetime.now().isoformat(),
            "metrics": metrics,
            "analysis": analysis,
            "doc_structure": state.get("doc_structure", {})
        }
        
        # Sauvegarde JSON final
        output_filename = f"{os.path.splitext(os.path.basename(state['file_path']))[0]}_analysis.json"
        output_path = os.path.join(OUTPUT_DIR, output_filename)
        
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(final_result, f, indent=2, ensure_ascii=False)
            
        print(f"✅ Analyse terminée: {output_path}")
        return {"metrics": metrics, "final_result": final_result}

    # ==================== HELPERS ====================

    def process_document(self, file_path: str, metadata: Dict = None) -> Dict[str, Any]:
        """Point d'entrée pour lancer le graphe
        
        Args:
            file_path: Path to the document to analyze
            metadata: Optional metadata dict from frontend questionnaire
        """
        print("\n" + "=" * 80)
        print(f"📊 TRAITEMENT: {os.path.basename(file_path)}")
        print("=" * 80)
        
        if metadata:
            print(f"📝 Métadonnées reçues du frontend: {metadata}")
        
        initial_state = {
            "file_path": file_path,
            "metadata_path": self.metadata_path,
            "metadata": metadata
        }
        
        result = self.workflow.invoke(initial_state)
        return result.get("final_result", {})

    def _summarize_document(self, doc_structure: Dict) -> str:
        """Prépare le contenu complet du document pour le prompt"""
        doc_type = doc_structure.get('type')
        content_text = []
        
        if doc_type == 'pptx':
            slides = doc_structure.get('slides', [])
            for slide in slides:
                slide_num = slide.get('slide_number')
                content_text.append(f"\n--- SLIDE {slide_num} ---")
                
                for item in slide.get('content', []):
                    item_id = item.get('id', 'N/A')
                    text = item.get('text', '')
                    if item.get('type') == 'table':
                        text = f"[TABLEAU] {json.dumps(item.get('data', []))}"
                    
                    content_text.append(f"ID: {item_id} | {text[:500]}") # Tronquer si trop long
                    
        elif doc_type == 'docx':
            content = doc_structure.get('content', [])
            for item in content:
                item_id = item.get('id', 'N/A')
                text = item.get('text', '')
                if item.get('type') == 'table':
                    text = f"[TABLEAU] {json.dumps(item.get('data', []))}"
                
                content_text.append(f"ID: {item_id} | {text[:500]}")
        
        elif doc_type == 'pdf':
            content = doc_structure.get('content', [])
            for item in content:
                page_num = item.get('page_number')
                text = item.get('text', '')
                content_text.append(f"\n--- PAGE {page_num} ---")
                content_text.append(f"{text[:2000]}") # Plus large pour PDF
        
        return "\n".join(content_text)

    def _parse_llm_response(self, response: str) -> Dict[str, Any]:
        """Parse la réponse JSON du LLM"""
        try:
            import re
            json_match = re.search(r'\{.*\}', response, re.DOTALL)
            
            if json_match:
                json_str = json_match.group(0)
                analysis = json.loads(json_str)
                return analysis
            else:
                print("  ⚠️ Pas de JSON valide dans la réponse LLM")
                return self._generate_fallback_analysis({})
        
        except Exception as e:
            print(f"  ❌ Erreur parsing réponse LLM: {e}")
            return self._generate_fallback_analysis({})

    def _generate_fallback_analysis(self, doc_structure: Dict) -> Dict[str, Any]:
        """Génère une analyse fallback si LLM échoue"""
        return {
            "conformite_globale": {
                "score": 0,
                "niveau": "inconnu"
            },
            "elements_conformes": [],
            "elements_non_conformes": [],
            "elements_manquants": [],
            "recommandations": ["Analyse LLM a échoué - vérification manuelle requise"],
            "fallback": True
        }

class ComplianceMetrics:
    """Calcule les métriques de conformité"""
    
    @staticmethod
    def calculate_metrics(analysis: Dict, expected_elements: int = None) -> Dict[str, float]:
        """Calcule les métriques de conformité"""
        
        conformes = len(analysis.get('elements_conformes', []))
        non_conformes = len(analysis.get('elements_non_conformes', []))
        manquants = len(analysis.get('elements_manquants', []))
        
        total_presents = conformes + non_conformes
        total_requis = total_presents + manquants
        
        metrics = {}
        
        if total_requis > 0:
            metrics['compliance_score'] = round(conformes / total_requis, 3)
            metrics['completeness'] = round(total_presents / total_requis, 3)
        else:
            metrics['compliance_score'] = 0.0
            metrics['completeness'] = 0.0
            
        if total_presents > 0:
            metrics['correctness'] = round(conformes / total_presents, 3)
        else:
            metrics['correctness'] = 0.0
            
        metrics['llm_score'] = analysis.get('conformite_globale', {}).get('score', 0) / 100.0
        
        metrics['total_conforme'] = conformes
        metrics['total_non_conforme'] = non_conformes
        metrics['total_manquant'] = manquants
        metrics['total_requis'] = total_requis
        
        return metrics

if __name__ == "__main__":
    # Test rapide
    agent = DocumentComplianceAgent()
    print("Agent LangGraph initialisé.")
